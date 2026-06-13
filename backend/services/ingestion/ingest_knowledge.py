import os
import sys
import json
from providers.llm_provider import get_embeddings_batch
from dotenv import load_dotenv
from pptx import Presentation
import concurrent.futures
from sqlalchemy import create_engine, text

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

load_dotenv()

# Database Configuration (Ensuring it uses +psycopg for SQLAlchemy)
DB_URL = os.getenv("DATABASE_URL", "postgresql+psycopg://root:root@localhost:5432/ai_db")
engine = create_engine(DB_URL, pool_pre_ping=True)

def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text("text") + "\n\n"
    return text

def extract_text_from_pptx(pptx_path):
    """Flat text extraction — kept as fallback. Prefer extract_slides_from_pptx."""
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n\n"
    return text

def extract_slides_from_pptx(pptx_path):
    """Structured per-slide extraction — one chunk per slide preserving semantic boundaries."""
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        slide_text = "\n".join(texts)
        if len(slide_text.strip()) >= 5:
            slides.append({
                "text": slide_text,
                "metadata": {"slide_number": i, "chunk_type": "slide"}
            })
    return slides

def chunk_text(text, chunk_size=1500, overlap=250):
    """Character-based chunking — kept as fallback. Prefer chunk_text_paragraphs."""
    chunks = []
    start = 0
    while start < len(text):
        chunks.append(text[start:start + chunk_size])
        start += chunk_size - overlap
    return chunks

def chunk_text_paragraphs(text, max_chars=1000, overlap=150):
    """Paragraph-aware PDF chunking: respects natural boundaries (double newlines)."""
    import re
    paragraphs = [p.strip() for p in re.split(r'\n\s*\n', text) if len(p.strip()) >= 5]
    chunks = []
    current = ""
    for para in paragraphs:
        if not current:
            current = para
        elif len(current) + len(para) + 2 <= max_chars:
            current += "\n\n" + para
        else:
            chunks.append(current)
            # Carry the tail of the previous chunk to preserve context continuity
            tail = current[-overlap:] if overlap > 0 else ""
            current = (tail + "\n\n" + para).strip() if tail else para
    if current.strip():
        chunks.append(current)
    return chunks

import time

def ingest_document(file_path, client_name="Internal", update_callback=None, brand_id=None, is_public=False, document_type="company_knowledge"):
    """
    STABLE INGESTION WITH BRAND SOVEREIGNTY, VISIBILITY & TAXONOMY (v12.0).
    """
    start_total = time.time()
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in [".pdf", ".pptx"]:
         raise ValueError(f"Format {ext} not supported.")
    
    source_filename = os.path.basename(file_path)
    print(f"\n[Worker] Starting ingestion for {client_name} (Brand ID: {brand_id}).", flush=True)
         
    if update_callback:
        update_callback(f"Reading document: {source_filename}", 5)
    
    start_extract = time.time()
    if ext == ".pdf":
        if not fitz: raise ValueError("PyMuPDF missing.")
        full_text = extract_text_from_pdf(file_path)
        raw_chunks = chunk_text_paragraphs(full_text)
        chunk_items = [{"text": c, "metadata": {"chunk_type": "paragraph"}} for c in raw_chunks]
    else:
        chunk_items = extract_slides_from_pptx(file_path)
    duration_extract = time.time() - start_extract

    valid_chunk_items = [c for c in chunk_items if len(c["text"].strip()) >= 5]
    total_chunks = len(valid_chunk_items)
    print(f"[Worker] Step 1: Extraction completed in {duration_extract:.2f}s ({total_chunks} semantic chunks).", flush=True)

    if total_chunks == 0:
        if update_callback:
            update_callback("No text found in document.", -1)
        return

    if update_callback:
        update_callback(f"Analyzing {total_chunks} semantic chunks...", 10)

    batch_size = 30
    inserted_total = 0
    start_sync = time.time()

    try:
        for i in range(0, total_chunks, batch_size):
            batch_items = valid_chunk_items[i:i+batch_size]
            batch_texts = [item["text"] for item in batch_items]
            current_count = i + len(batch_items)
            current_batch_idx = (i // batch_size) + 1

            if update_callback:
                perc = 10 + int((i / total_chunks) * 85)
                update_callback(f"Indexing fragments ({current_count}/{total_chunks})...", perc)

            try:
                batch_embeddings = get_embeddings_batch(batch_texts)
                if not batch_embeddings or len(batch_embeddings) == 0:
                    print(f"  [Worker] Batch {current_batch_idx} skipped: No valid content or embeddings found.", flush=True)
                    continue

                # Each batch gets a fresh connection+transaction — a failed batch never
                # corrupts the connection state for subsequent batches (psycopg3 + SA 2.x).
                with engine.begin() as conn:
                    for item, emb in zip(batch_items, batch_embeddings):
                        if emb is None: continue
                        emb_pg = f"[{','.join(map(str, emb))}]"
                        chunk_meta = item.get("metadata", {})
                        conn.execute(
                            text("""
                            INSERT INTO corporate_knowledge (content, meta_data, embedding, brand_id, is_public, source_filename, document_type)
                            VALUES (:content, :meta_data, cast(:embedding as vector), :brand_id, :is_public, :source_filename, :document_type)
                            """),
                            {
                                "content": item["text"],
                                "meta_data": json.dumps({
                                    "source": source_filename,
                                    "client": client_name,
                                    "brand_id": brand_id,
                                    "is_public": is_public,
                                    "document_type": document_type,
                                    **chunk_meta
                                }),
                                "embedding": emb_pg,
                                "brand_id": brand_id,
                                "is_public": 1 if is_public else 0,
                                "source_filename": source_filename,
                                "document_type": document_type
                            }
                        )
                        inserted_total += 1
            except Exception as batch_err:
                print(f"  [Error] Batch {current_batch_idx} skipped due to error: {batch_err}", flush=True)
                continue
    except Exception as e:
        print(f"  [Fatal] Knowledge Sync Error: {str(e)}", flush=True)
        if update_callback:
            update_callback(f"Fatal Sync Error: {str(e)}", -1)
        raise e

    duration_sync = time.time() - start_sync
    duration_total = time.time() - start_total
    
    if update_callback:
        update_callback(f"Synchronization finalized ({inserted_total} items).", 100)
    
    print(f"[Worker] Step 2: RAG sync completed in {duration_sync:.2f}s ({inserted_total} records).", flush=True)
    print(f"[Worker] Total Ingestion Time: {duration_total:.2f}s\n", flush=True)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python ingest_knowledge.py <file_path> [client_name]")
        sys.exit(1)
    
    p = sys.argv[1]
    c = sys.argv[2] if len(sys.argv) > 2 else "Internal"
    if os.path.exists(p):
        ingest_document(p, c)
