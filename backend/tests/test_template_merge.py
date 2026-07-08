"""
test_template_merge.py — Unit tests for the Template Merge Engine (v2).

Covers the pure/mockable logic in services/templates/*: shared traversal
(groups, tables, depth cap), role inference, char-limit estimation, action
classification (analyzer), LLM response unwrapping and markdown stripping
(content), in-place text replacement preserving run formatting including
bullet-aware paragraph mapping and the empty-slot policy (renderer), the
merge report, and the orchestrator's job lifecycle and error handling.
No DB or real LLM calls — DB and generate_json/search_rag are mocked so
these run without the test container.
"""
from types import SimpleNamespace
from unittest.mock import MagicMock, patch

import pytest

from services.templates.template_config import TemplateMergeConfig
from services.templates import template_analyzer as analyzer_mod
from services.templates import template_content as content_mod
from services.templates import template_merge_orchestrator as orch_mod
from services.templates import template_renderer as renderer_mod
from services.templates.template_traversal import TextTarget, collect_text_targets


def make_config(**overrides) -> TemplateMergeConfig:
    return TemplateMergeConfig(**overrides)


def make_target(width=0, height=0, top=0, is_placeholder=False,
                placeholder_type=None, kind="shape", key="1", name="t"):
    shape = None
    if is_placeholder:
        shape = SimpleNamespace(placeholder_format=SimpleNamespace(type=placeholder_type))
    return TextTarget(
        key=key, text_frame=None, name=name, kind=kind,
        is_placeholder=is_placeholder, shape=shape,
        width=width, height=height, top=top,
    )


# ---------------------------------------------------------------------------
# template_traversal — collect_text_targets (real python-pptx objects)
# ---------------------------------------------------------------------------

def _blank_slide():
    from pptx import Presentation
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


@pytest.mark.unit
def test_collect_targets_plain_textbox():
    from pptx.util import Inches
    _, slide = _blank_slide()
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    box.text_frame.text = "hello"

    targets, preserved = collect_text_targets(slide, max_group_depth=3)

    assert len(targets) == 1
    assert targets[0].key == str(box.shape_id)
    assert targets[0].kind == "shape"
    assert preserved == 0


@pytest.mark.unit
def test_collect_targets_recurses_into_groups_with_path_keys():
    from pptx.util import Inches
    _, slide = _blank_slide()
    tb1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tb1.text_frame.text = "inside one"
    tb2 = slide.shapes.add_textbox(Inches(3), Inches(1), Inches(2), Inches(1))
    tb2.text_frame.text = "inside two"
    group = slide.shapes.add_group_shape([tb1, tb2])

    targets, _ = collect_text_targets(slide, max_group_depth=3)

    keys = sorted(t.key for t in targets)
    assert len(targets) == 2
    assert all(k.startswith(f"{group.shape_id}/") for k in keys)
    assert all(t.kind == "group_child" for t in targets)


@pytest.mark.unit
def test_collect_targets_group_beyond_depth_cap_is_preserved():
    from pptx.util import Inches
    _, slide = _blank_slide()
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tb.text_frame.text = "unreachable"
    slide.shapes.add_group_shape([tb])

    targets, preserved = collect_text_targets(slide, max_group_depth=0)

    assert targets == []
    assert preserved == 1


@pytest.mark.unit
def test_collect_targets_table_cells_with_rc_keys_and_spanned_skipped():
    from pptx.util import Inches
    _, slide = _blank_slide()
    frame = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    table = frame.table
    for r in range(2):
        for c in range(2):
            table.cell(r, c).text = f"cell {r}{c}"
    table.cell(0, 0).merge(table.cell(0, 1))  # (0,1) becomes spanned

    targets, preserved = collect_text_targets(slide, max_group_depth=3)

    keys = sorted(t.key for t in targets)
    sid = frame.shape_id
    assert keys == [f"{sid}:r0c0", f"{sid}:r1c0", f"{sid}:r1c1"]
    assert all(t.kind == "cell" for t in targets)
    assert preserved == 1  # the spanned cell


@pytest.mark.unit
def test_collect_targets_cell_geometry_from_column_and_row():
    from pptx.util import Inches
    _, slide = _blank_slide()
    frame = slide.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(4), Inches(1))
    frame.table.cell(0, 0).text = "x"

    targets, _ = collect_text_targets(slide, max_group_depth=3)

    cell0 = next(t for t in targets if t.key.endswith(":r0c0"))
    assert cell0.width == frame.table.columns[0].width
    assert cell0.height == frame.table.rows[0].height


# ---------------------------------------------------------------------------
# template_analyzer — _area_within_bounds
# ---------------------------------------------------------------------------

@pytest.mark.unit
def test_area_within_bounds_accepts_normal_shape():
    config = make_config()
    target = make_target(width=100, height=100)
    assert analyzer_mod._area_within_bounds(target, slide_area_emu=100_000, config=config) is True


@pytest.mark.unit
def test_area_within_bounds_rejects_background():
    config = make_config(shape_bg_area_threshold=0.80)
    target = make_target(width=1000, height=1000)  # ratio = 1.0 > 0.80
    assert analyzer_mod._area_within_bounds(target, slide_area_emu=1_000_000, config=config) is False


@pytest.mark.unit
def test_area_within_bounds_rejects_decorative_dot():
    config = make_config(shape_min_area_threshold=0.005)
    target = make_target(width=10, height=10)  # ratio tiny
    assert analyzer_mod._area_within_bounds(target, slide_area_emu=1_000_000, config=config) is False


@pytest.mark.unit
def test_area_within_bounds_zero_slide_area_defaults_true():
    config = make_config()
    target = make_target(width=10, height=10)
    assert analyzer_mod._area_within_bounds(target, slide_area_emu=0, config=config) is True


# ---------------------------------------------------------------------------
# template_analyzer — _infer_role
# ---------------------------------------------------------------------------

@pytest.mark.unit
def test_infer_role_placeholder_title():
    from pptx.enum.shapes import PP_PLACEHOLDER
    config = make_config()
    target = make_target(is_placeholder=True, placeholder_type=PP_PLACEHOLDER.TITLE)
    assert analyzer_mod._infer_role(target, slide_height_emu=1000, slide_area_emu=1_000_000, config=config) == "title"


@pytest.mark.unit
def test_infer_role_placeholder_body():
    from pptx.enum.shapes import PP_PLACEHOLDER
    config = make_config()
    target = make_target(is_placeholder=True, placeholder_type=PP_PLACEHOLDER.BODY)
    assert analyzer_mod._infer_role(target, slide_height_emu=1000, slide_area_emu=1_000_000, config=config) == "body"


@pytest.mark.unit
def test_infer_role_cell_is_always_body():
    config = make_config(footnote_area_fraction=0.03)
    # tiny area would classify a shape as footnote — cells must stay body
    target = make_target(width=10, height=10, top=5000, kind="cell")
    assert analyzer_mod._infer_role(target, slide_height_emu=10000, slide_area_emu=1_000_000, config=config) == "body"


@pytest.mark.unit
def test_infer_role_non_placeholder_footnote_by_small_area():
    config = make_config(footnote_area_fraction=0.03)
    target = make_target(width=10, height=10, top=5000)
    assert analyzer_mod._infer_role(target, slide_height_emu=10000, slide_area_emu=1_000_000, config=config) == "footnote"


@pytest.mark.unit
def test_infer_role_non_placeholder_title_by_top_position():
    config = make_config(footnote_area_fraction=0.001, title_top_fraction=0.20)
    target = make_target(width=500, height=500, top=100)
    assert analyzer_mod._infer_role(target, slide_height_emu=10000, slide_area_emu=1_000_000, config=config) == "title"


@pytest.mark.unit
def test_infer_role_non_placeholder_defaults_body():
    config = make_config(footnote_area_fraction=0.001, title_top_fraction=0.05)
    target = make_target(width=500, height=500, top=9000)
    assert analyzer_mod._infer_role(target, slide_height_emu=10000, slide_area_emu=1_000_000, config=config) == "body"


# ---------------------------------------------------------------------------
# template_analyzer — _estimate_char_limit
# ---------------------------------------------------------------------------

@pytest.mark.unit
def test_estimate_char_limit_title_short_hint_uses_multiplier():
    config = make_config(short_hint_threshold=15, short_hint_title_multiplier=3)
    target = make_target()
    assert analyzer_mod._estimate_char_limit(target, "title", "$45", config) == max(len("$45") * 3, 20)


@pytest.mark.unit
def test_estimate_char_limit_title_long_hint_uses_default():
    config = make_config(title_char_limit=80)
    target = make_target()
    hint = "A" * 30  # longer than short_hint_threshold (15)
    assert analyzer_mod._estimate_char_limit(target, "title", hint, config) == 80


@pytest.mark.unit
def test_estimate_char_limit_footnote_uses_fixed_limit():
    config = make_config(footnote_char_limit=120)
    target = make_target()
    assert analyzer_mod._estimate_char_limit(target, "footnote", "any hint", config) == 120


@pytest.mark.unit
def test_estimate_char_limit_body_short_hint_uses_multiplier():
    config = make_config(short_hint_threshold=15, short_hint_body_multiplier=4)
    target = make_target()
    assert analyzer_mod._estimate_char_limit(target, "body", "23%", config) == max(len("23%") * 4, 30)


@pytest.mark.unit
def test_estimate_char_limit_body_area_based_clamped_to_min():
    config = make_config(body_char_limit_min=80, body_char_limit_max=600, chars_per_sq_inch=30)
    target = make_target(width=91440, height=91440)  # 0.1in × 0.1in
    hint = "A" * 30
    assert analyzer_mod._estimate_char_limit(target, "body", hint, config) == 80


@pytest.mark.unit
def test_estimate_char_limit_body_area_based_clamped_to_max():
    config = make_config(body_char_limit_min=80, body_char_limit_max=600, chars_per_sq_inch=30)
    target = make_target(width=914400 * 20, height=914400 * 20)  # huge box
    hint = "A" * 30
    assert analyzer_mod._estimate_char_limit(target, "body", hint, config) == 600


# ---------------------------------------------------------------------------
# template_analyzer — _infer_action
# ---------------------------------------------------------------------------

@pytest.mark.unit
def test_infer_action_footnote_always_preserved():
    config = make_config()
    assert analyzer_mod._infer_action(is_placeholder=False, role="footnote", hint="anything", config=config) == "preserve"


@pytest.mark.unit
def test_infer_action_preserve_keyword_overrides_length():
    config = make_config(preserve_keywords="confidential,proprietary")
    long_hint = "This document is Confidential " + "x" * 200
    assert analyzer_mod._infer_action(is_placeholder=False, role="body", hint=long_hint, config=config) == "preserve"


@pytest.mark.unit
def test_infer_action_placeholder_always_rewrite():
    config = make_config()
    assert analyzer_mod._infer_action(is_placeholder=True, role="body", hint="", config=config) == "rewrite"


@pytest.mark.unit
def test_infer_action_short_hint_preserved():
    config = make_config(preserve_max_hint_chars=50)
    assert analyzer_mod._infer_action(is_placeholder=False, role="body", hint="Short label", config=config) == "preserve"


@pytest.mark.unit
def test_infer_action_medium_hint_adapt():
    config = make_config(preserve_max_hint_chars=10, adapt_max_hint_chars=150)
    hint = "A" * 100
    assert analyzer_mod._infer_action(is_placeholder=False, role="body", hint=hint, config=config) == "adapt"


@pytest.mark.unit
def test_infer_action_long_hint_rewrite():
    config = make_config(preserve_max_hint_chars=10, adapt_max_hint_chars=50)
    hint = "A" * 200
    assert analyzer_mod._infer_action(is_placeholder=False, role="body", hint=hint, config=config) == "rewrite"


# ---------------------------------------------------------------------------
# template_analyzer — analyze_template over groups/tables (real pptx)
# ---------------------------------------------------------------------------

@pytest.mark.unit
def test_analyze_template_covers_group_children_and_cells(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tb.text_frame.text = "Grouped text long enough to be considered adaptable by the classifier here"
    slide.shapes.add_group_shape([tb])
    frame = slide.shapes.add_table(1, 1, Inches(1), Inches(3), Inches(3), Inches(1))
    frame.table.cell(0, 0).text = "Cell text long enough to be adapted rather than preserved by length rules"
    path = str(tmp_path / "t.pptx")
    prs.save(path)

    profiles = analyzer_mod.analyze_template(path, make_config())

    kinds = {s.kind for s in profiles[0].slots}
    assert "group_child" in kinds
    assert "cell" in kinds
    assert all(isinstance(s.slot_key, str) for s in profiles[0].slots)


# ---------------------------------------------------------------------------
# template_content — _unwrap_value
# ---------------------------------------------------------------------------

@pytest.mark.unit
def test_unwrap_value_plain_string_passthrough():
    assert content_mod._unwrap_value("Hello world", max_bullet_items=6) == "Hello world"


@pytest.mark.unit
def test_unwrap_value_dict_with_content_key():
    value = {"role": "body", "content": "Actual text"}
    assert content_mod._unwrap_value(value, max_bullet_items=6) == "Actual text"


@pytest.mark.unit
def test_unwrap_value_dict_without_known_key_joins_values():
    value = {"foo": "bar", "baz": "qux"}
    result = content_mod._unwrap_value(value, max_bullet_items=6)
    assert "bar" in result and "qux" in result


@pytest.mark.unit
def test_unwrap_value_list_joins_with_newline_and_caps_items():
    value = ["one", "two", "three", "four"]
    result = content_mod._unwrap_value(value, max_bullet_items=2)
    assert result == "one\ntwo"


@pytest.mark.unit
def test_unwrap_value_string_repr_of_dict_is_parsed():
    value = "{'role': 'body', 'content': 'Parsed text'}"
    assert content_mod._unwrap_value(value, max_bullet_items=6) == "Parsed text"


@pytest.mark.unit
def test_unwrap_value_malformed_string_repr_falls_back_to_text():
    value = "{'content': unterminated"
    result = content_mod._unwrap_value(value, max_bullet_items=6)
    assert result == value


# ---------------------------------------------------------------------------
# template_content / template_renderer — _strip_markdown
# ---------------------------------------------------------------------------

@pytest.mark.unit
@pytest.mark.parametrize("mod", [content_mod, renderer_mod])
def test_strip_markdown_removes_bold_italic_headings_bullets(mod):
    assert mod._strip_markdown("**Bold** and *italic*") == "Bold and italic"
    assert mod._strip_markdown("## Heading") == "Heading"
    assert mod._strip_markdown("- bullet item") == "bullet item"
    assert mod._strip_markdown("`code`") == "code"


@pytest.mark.unit
def test_strip_markdown_empty_string():
    assert content_mod._strip_markdown("") == ""
    assert content_mod._strip_markdown(None) is None


# ---------------------------------------------------------------------------
# template_content — _generate_for_slide (mocked LLM + RAG)
# ---------------------------------------------------------------------------

def make_slot(slot_key="1", role="body", action="rewrite", char_limit=80, hint="hint", kind="shape"):
    return analyzer_mod.TextSlot(
        slide_idx=0, slot_key=slot_key, shape_name=f"shape{slot_key}", role=role,
        char_limit=char_limit, hint=hint, is_placeholder=False, action=action, kind=kind,
    )


def make_profile(slots):
    profile = analyzer_mod.SlideProfile(slide_idx=0, slide_width_emu=1, slide_height_emu=1)
    profile.slots = slots
    return profile


@pytest.mark.unit
def test_generate_for_slide_skips_llm_when_all_preserved():
    config = make_config()
    profile = make_profile([make_slot(action="preserve")])
    with patch("services.templates.template_content.generate_json") as mock_llm:
        result = content_mod._generate_for_slide(
            profile=profile, knowledge_filename="doc.pdf", brand_id=1,
            user_prompt="prompt", slide_num=1, total_slides=1, config=config,
        )
    assert result == {}
    mock_llm.assert_not_called()


@pytest.mark.unit
def test_generate_for_slide_maps_slot_keys_to_llm_content():
    config = make_config()
    slot = make_slot(slot_key="42:r1c0", char_limit=100, kind="cell")
    profile = make_profile([slot])
    with patch("services.templates.template_content.generate_json", return_value={"42:r1c0": "Generated text"}), \
         patch("services.templates.template_content.search_rag", return_value=["context chunk"]):
        result = content_mod._generate_for_slide(
            profile=profile, knowledge_filename="doc.pdf", brand_id=1,
            user_prompt="prompt", slide_num=1, total_slides=1, config=config,
        )
    assert result == {"42:r1c0": "Generated text"}


@pytest.mark.unit
def test_generate_for_slide_truncates_over_char_limit():
    config = make_config()
    slot = make_slot(slot_key="1", char_limit=10)
    profile = make_profile([slot])
    long_text = "This is a very long sentence that exceeds the limit"
    with patch("services.templates.template_content.generate_json", return_value={"1": long_text}), \
         patch("services.templates.template_content.search_rag", return_value=[]):
        result = content_mod._generate_for_slide(
            profile=profile, knowledge_filename="doc.pdf", brand_id=1,
            user_prompt="prompt", slide_num=1, total_slides=1, config=config,
        )
    assert len(result["1"]) <= 11  # char_limit + ellipsis
    assert result["1"].endswith("…")


@pytest.mark.unit
def test_generate_for_slide_rag_failure_still_calls_llm():
    config = make_config()
    slot = make_slot(slot_key="1", char_limit=100)
    profile = make_profile([slot])
    with patch("services.templates.template_content.generate_json", return_value={"1": "ok"}) as mock_llm, \
         patch("services.templates.template_content.search_rag", side_effect=RuntimeError("rag down")):
        result = content_mod._generate_for_slide(
            profile=profile, knowledge_filename="doc.pdf", brand_id=1,
            user_prompt="prompt", slide_num=1, total_slides=1, config=config,
        )
    mock_llm.assert_called_once()
    assert result == {"1": "ok"}


@pytest.mark.unit
def test_generate_slide_contents_empty_slots_slide_skips_llm():
    config = make_config()
    empty_profile = make_profile([])
    with patch("services.templates.template_content.generate_json") as mock_llm:
        results = content_mod.generate_slide_contents(
            profiles=[empty_profile], knowledge_filename="doc.pdf", brand_id=1,
            user_prompt="prompt", config=config,
        )
    assert results == [{}]
    mock_llm.assert_not_called()


@pytest.mark.unit
def test_generate_slide_contents_per_slide_error_yields_none():
    config = make_config()
    slot = make_slot(slot_key="7")
    profile = make_profile([slot])
    with patch("services.templates.template_content._generate_for_slide", side_effect=RuntimeError("boom")):
        results = content_mod.generate_slide_contents(
            profiles=[profile], knowledge_filename="doc.pdf", brand_id=1,
            user_prompt="prompt", config=config,
        )
    # A failed slide yields None so the renderer keeps the original text
    # (reported as `failed`) instead of blanking it.
    assert results == [None]


# ---------------------------------------------------------------------------
# template_plan — plan_deck (v2 Fase 2, mocked LLM + RAG)
# ---------------------------------------------------------------------------

from services.templates import template_plan as plan_mod
from services.templates.template_plan import DeckPlan, SlidePlan


def _two_profiles():
    p0 = make_profile([make_slot(slot_key="1", action="rewrite", hint="Old title one")])
    p1 = analyzer_mod.SlideProfile(slide_idx=1, slide_width_emu=1, slide_height_emu=1)
    p1.slots = [make_slot(slot_key="2", action="adapt", hint="Old metric text")]
    return [p0, p1]


def _valid_plan_json():
    return {
        "language": "es",
        "tone": "profesional y directo",
        "slides": [
            {"slide": 1, "topic": "Apertura", "key_points": ["quiénes somos"], "rag_query": "empresa historia equipo"},
            {"slide": 2, "topic": "Métricas", "key_points": ["ingresos", "clientes"], "rag_query": "resultados financieros 2025"},
        ],
    }


@pytest.mark.unit
def test_plan_deck_disabled_skips_llm():
    config = make_config(outline_enabled=False)
    with patch("services.templates.template_plan.generate_json") as mock_llm:
        assert plan_mod.plan_deck(_two_profiles(), "doc.pdf", 1, "prompt", config) is None
    mock_llm.assert_not_called()


@pytest.mark.unit
def test_plan_deck_no_active_slots_skips_llm():
    config = make_config()
    profile = make_profile([make_slot(action="preserve")])
    with patch("services.templates.template_plan.generate_json") as mock_llm:
        assert plan_mod.plan_deck([profile], "doc.pdf", 1, "prompt", config) is None
    mock_llm.assert_not_called()


@pytest.mark.unit
def test_plan_deck_happy_path_maps_by_position_onto_slide_idx():
    config = make_config()
    with patch("services.templates.template_plan.generate_json", return_value=_valid_plan_json()), \
         patch("services.templates.template_plan.search_rag", return_value="chunk one\n---\nchunk two"):
        plan = plan_mod.plan_deck(_two_profiles(), "doc.pdf", 1, "prompt", config)

    assert isinstance(plan, DeckPlan)
    assert plan.language == "es"
    assert plan.for_slide(0).topic == "Apertura"
    assert plan.for_slide(1).rag_query == "resultados financieros 2025"


@pytest.mark.unit
def test_plan_deck_llm_failure_degrades_to_none():
    config = make_config()
    with patch("services.templates.template_plan.generate_json", side_effect=RuntimeError("llm down")), \
         patch("services.templates.template_plan.search_rag", return_value=""):
        assert plan_mod.plan_deck(_two_profiles(), "doc.pdf", 1, "prompt", config) is None


@pytest.mark.unit
def test_plan_deck_rag_failure_still_calls_llm():
    config = make_config()
    with patch("services.templates.template_plan.generate_json", return_value=_valid_plan_json()) as mock_llm, \
         patch("services.templates.template_plan.search_rag", side_effect=RuntimeError("rag down")):
        plan = plan_mod.plan_deck(_two_profiles(), "doc.pdf", 1, "prompt", config)
    mock_llm.assert_called_once()
    assert plan is not None


@pytest.mark.unit
def test_parse_plan_malformed_shapes_return_none():
    profiles = _two_profiles()
    assert plan_mod._parse_plan("not a dict", profiles) is None
    assert plan_mod._parse_plan({}, profiles) is None
    assert plan_mod._parse_plan({"slides": "nope"}, profiles) is None
    assert plan_mod._parse_plan({"slides": []}, profiles) is None


@pytest.mark.unit
def test_parse_plan_invalid_entry_dropped_valid_ones_kept():
    profiles = _two_profiles()
    raw = {
        "language": "EN ",
        "tone": "crisp",
        "slides": [
            {"slide": 1, "topic": "", "key_points": [], "rag_query": "x"},          # invalid: empty topic
            {"slide": 2, "topic": "Metrics", "key_points": "oops", "rag_query": "q"},  # tolerated: key_points not a list
        ],
    }
    plan = plan_mod._parse_plan(raw, profiles)
    assert plan is not None
    assert plan.language == "en"
    assert plan.for_slide(0) is None          # slide 1 degraded to v1
    assert plan.for_slide(1).topic == "Metrics"
    assert plan.for_slide(1).key_points == []


# ---------------------------------------------------------------------------
# template_content — v2 Fase 2 (plan integration, dedup, language, summaries)
# ---------------------------------------------------------------------------

@pytest.mark.unit
def test_generate_for_slide_uses_plan_rag_query():
    config = make_config()
    profile = make_profile([make_slot(slot_key="1", hint="Old template hint")])
    plan_slide = SlidePlan(topic="New topic", key_points=["a"], rag_query="planned query")
    with patch("services.templates.template_content.generate_json", return_value={"1": "ok"}), \
         patch("services.templates.template_content.search_rag", return_value="chunk") as mock_rag:
        content_mod._generate_for_slide(
            profile=profile, knowledge_filename="doc.pdf", brand_id=1,
            user_prompt="prompt", slide_num=1, total_slides=1, config=config,
            plan_slide=plan_slide,
        )
    assert mock_rag.call_args.kwargs["query"] == "planned query"


@pytest.mark.unit
def test_generate_for_slide_without_plan_keeps_v1_hint_query():
    config = make_config()
    profile = make_profile([make_slot(slot_key="1", role="title", hint="Old title")])
    with patch("services.templates.template_content.generate_json", return_value={"1": "ok"}), \
         patch("services.templates.template_content.search_rag", return_value="chunk") as mock_rag:
        content_mod._generate_for_slide(
            profile=profile, knowledge_filename="doc.pdf", brand_id=1,
            user_prompt="prompt", slide_num=1, total_slides=1, config=config,
        )
    assert mock_rag.call_args.kwargs["query"] == "Old title"


@pytest.mark.unit
def test_build_prompt_includes_plan_prev_and_language_sections():
    plan_slide = SlidePlan(topic="Growth story", key_points=["12% up"], rag_query="q")
    prompt = content_mod._build_prompt(
        slide_num=2, total_slides=3, topic_hint="hint", user_prompt="intent",
        rag_context="ctx", slots_desc="slots",
        plan_slide=plan_slide, language="es", tone="direct",
        prev_summaries=["Slide 1: Opening about the company"],
    )
    assert "This slide's topic: Growth story" in prompt
    assert "- 12% up" in prompt
    assert "Deck tone: direct" in prompt
    assert "Slide 1: Opening about the company" in prompt
    assert 'write ALL content in "es"' in prompt


@pytest.mark.unit
def test_build_prompt_without_plan_falls_back_to_user_intent_language():
    prompt = content_mod._build_prompt(
        slide_num=1, total_slides=1, topic_hint="hint", user_prompt="intent",
        rag_context="ctx", slots_desc="slots",
    )
    assert "DECK PLAN" not in prompt
    assert "PREVIOUS SLIDES" not in prompt
    assert "same language as the USER INTENT" in prompt


@pytest.mark.unit
def test_deprioritize_used_chunks_fresh_first_and_stale_dropped_by_cap():
    used = {"chunk B"}
    ctx = "chunk A\n---\nchunk B\n---\nchunk C"
    result = content_mod._deprioritize_used_chunks(ctx, used, max_chars=20)
    # 20 chars fit both fresh chunks (7 + 5 sep + 7 = 19); the stale B is
    # pushed last and dropped by the cap
    assert result == "chunk A\n---\nchunk C"
    assert "chunk A" in used and "chunk C" in used


@pytest.mark.unit
def test_deprioritize_used_chunks_stale_kept_when_budget_allows():
    used = {"chunk B"}
    ctx = "chunk A\n---\nchunk B"
    result = content_mod._deprioritize_used_chunks(ctx, used, max_chars=200)
    assert result == "chunk A\n---\nchunk B"  # fresh first, stale after


@pytest.mark.unit
def test_summarize_slide_prefers_title_and_truncates():
    title_slot = make_slot(slot_key="t", role="title")
    body_slot = make_slot(slot_key="b", role="body")
    profile = make_profile([body_slot, title_slot])
    content = {"b": "body text first in dict", "t": "T" * 300}
    summary = content_mod._summarize_slide(profile, content)
    assert summary == "T" * 100


@pytest.mark.unit
def test_generate_slide_contents_accumulates_prev_summaries():
    config = make_config()
    profiles = _two_profiles()
    captured_prompts = []

    def fake_generate(prompt):
        captured_prompts.append(prompt)
        return {"1": "First slide generated title", "2": "Second slide text"}

    with patch("services.templates.template_content.generate_json", side_effect=fake_generate), \
         patch("services.templates.template_content.search_rag", return_value=""):
        results = content_mod.generate_slide_contents(
            profiles=profiles, knowledge_filename="doc.pdf", brand_id=1,
            user_prompt="prompt", config=config,
        )

    assert len(results) == 2
    assert "PREVIOUS SLIDES" not in captured_prompts[0]
    assert "Slide 1: First slide generated title" in captured_prompts[1]


# ---------------------------------------------------------------------------
# template_renderer — _replace_text_frame / bullets / blank (real pptx objects)
# ---------------------------------------------------------------------------

@pytest.mark.unit
def test_replace_text_frame_preserves_formatting_and_sets_text():
    from pptx.util import Inches, Pt
    _, slide = _blank_slide()
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tf = box.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "Original"
    run.font.bold = True
    run.font.size = Pt(24)

    renderer_mod._replace_text_frame(tf, "Replaced text")

    assert tf.paragraphs[0].runs[0].text == "Replaced text"
    assert tf.paragraphs[0].runs[0].font.bold is True
    assert tf.paragraphs[0].runs[0].font.size == Pt(24)


@pytest.mark.unit
def test_replace_text_frame_multiline_single_paragraph_uses_soft_breaks():
    from pptx.util import Inches
    _, slide = _blank_slide()
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
    tf = box.text_frame
    tf.paragraphs[0].add_run().text = "Original"

    renderer_mod._replace_text_frame(tf, "Line one\nLine two")

    # Single original paragraph → v1 behavior: both lines in paragraph 0
    full_text = tf.paragraphs[0].text
    assert "Line one" in full_text and "Line two" in full_text
    assert len(tf.paragraphs) == 1


@pytest.mark.unit
def test_replace_text_frame_bulleted_maps_line_per_paragraph():
    from pptx.util import Inches, Pt
    _, slide = _blank_slide()
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(3))
    tf = box.text_frame
    r0 = tf.paragraphs[0].add_run(); r0.text = "Old bullet one"; r0.font.size = Pt(18)
    p1 = tf.add_paragraph(); r1 = p1.add_run(); r1.text = "Old bullet two"; r1.font.size = Pt(12)
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = "Old bullet three"; r2.font.size = Pt(12)

    renderer_mod._replace_text_frame(tf, "New one\nNew two")

    texts = [p.text for p in tf.paragraphs]
    assert texts[0] == "New one"
    assert texts[1] == "New two"
    assert texts[2] == ""  # leftover paragraph blanked, geometry stable
    # each line kept ITS paragraph's formatting, not the frame-wide first run's
    assert tf.paragraphs[0].runs[0].font.size == Pt(18)
    assert tf.paragraphs[1].runs[0].font.size == Pt(12)


@pytest.mark.unit
def test_replace_text_frame_bulleted_more_lines_than_paragraphs_clones_last():
    from pptx.util import Inches
    _, slide = _blank_slide()
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(3))
    tf = box.text_frame
    tf.paragraphs[0].add_run().text = "Old one"
    tf.add_paragraph().add_run().text = "Old two"

    renderer_mod._replace_text_frame(tf, "A\nB\nC\nD")

    texts = [p.text for p in tf.paragraphs if p.text]
    assert texts == ["A", "B", "C", "D"]


@pytest.mark.unit
def test_blank_text_frame_clears_all_text():
    from pptx.util import Inches
    _, slide = _blank_slide()
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tf = box.text_frame
    tf.paragraphs[0].add_run().text = "Old"
    tf.add_paragraph().add_run().text = "Older"

    renderer_mod._blank_text_frame(tf)

    assert tf.text.strip() == ""


# ---------------------------------------------------------------------------
# template_renderer — _merge_slide outcomes + empty policy (real pptx objects)
# ---------------------------------------------------------------------------

def _slide_with_two_boxes():
    from pptx.util import Inches
    prs, slide = _blank_slide()
    box1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    box1.text_frame.paragraphs[0].add_run().text = "Original 1"
    box2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
    box2.text_frame.paragraphs[0].add_run().text = "Original 2"
    return slide, str(box1.shape_id), str(box2.shape_id)


def _profile_for(slide_idx, slots):
    prof = analyzer_mod.SlideProfile(slide_idx=slide_idx, slide_width_emu=1, slide_height_emu=1)
    prof.slots = slots
    return prof


@pytest.mark.unit
def test_merge_slide_outcomes_rewritten_preserved_and_untouched_shapes():
    slide, k1, k2 = _slide_with_two_boxes()
    profile = _profile_for(0, [
        make_slot(slot_key=k1, action="rewrite"),
        make_slot(slot_key=k2, action="preserve"),
    ])
    entries = renderer_mod._merge_slide(slide, profile, {k1: "New text"}, make_config())

    by_key = {e["key"]: e["outcome"] for e in entries}
    assert by_key == {k1: "rewritten", k2: "preserved"}


@pytest.mark.unit
def test_merge_slide_empty_rewrite_blank_policy_blanks_and_reports_unfilled():
    slide, k1, _ = _slide_with_two_boxes()
    profile = _profile_for(0, [make_slot(slot_key=k1, action="rewrite")])

    entries = renderer_mod._merge_slide(
        slide, profile, {k1: ""}, make_config(empty_rewrite_policy="blank"))

    assert entries[0]["outcome"] == "unfilled"
    target = next(t for t in collect_text_targets(slide, 3)[0] if t.key == k1)
    assert target.text_frame.text.strip() == ""


@pytest.mark.unit
def test_merge_slide_empty_rewrite_keep_policy_keeps_original():
    slide, k1, _ = _slide_with_two_boxes()
    profile = _profile_for(0, [make_slot(slot_key=k1, action="rewrite")])

    entries = renderer_mod._merge_slide(
        slide, profile, {k1: ""}, make_config(empty_rewrite_policy="keep"))

    assert entries[0]["outcome"] == "kept_original"
    target = next(t for t in collect_text_targets(slide, 3)[0] if t.key == k1)
    assert target.text_frame.text == "Original 1"


@pytest.mark.unit
def test_merge_slide_empty_adapt_keeps_original_even_with_blank_policy():
    slide, k1, _ = _slide_with_two_boxes()
    profile = _profile_for(0, [make_slot(slot_key=k1, action="adapt")])

    entries = renderer_mod._merge_slide(
        slide, profile, {k1: ""}, make_config(empty_rewrite_policy="blank"))

    assert entries[0]["outcome"] == "kept_original"


@pytest.mark.unit
def test_merge_slide_none_content_map_reports_failed_and_keeps_text():
    slide, k1, _ = _slide_with_two_boxes()
    profile = _profile_for(0, [make_slot(slot_key=k1, action="rewrite")])

    entries = renderer_mod._merge_slide(slide, profile, None, make_config())

    assert entries[0]["outcome"] == "failed"
    target = next(t for t in collect_text_targets(slide, 3)[0] if t.key == k1)
    assert target.text_frame.text == "Original 1"


@pytest.mark.unit
def test_merge_slide_unresolvable_slot_key_reports_failed():
    slide, k1, _ = _slide_with_two_boxes()
    profile = _profile_for(0, [make_slot(slot_key="9999", action="rewrite")])

    entries = renderer_mod._merge_slide(slide, profile, {"9999": "text"}, make_config())

    assert entries[0]["outcome"] == "failed"


@pytest.mark.unit
def test_summarize_counts_outcomes():
    report_slides = [
        {"slide": 0, "slots": [{"outcome": "rewritten"}, {"outcome": "preserved"}], "preserved_shapes": 2},
        {"slide": 1, "slots": [{"outcome": "rewritten"}, {"outcome": "unfilled"}], "preserved_shapes": 0},
    ]
    summary = renderer_mod._summarize(report_slides)
    assert summary["rewritten"] == 2
    assert summary["preserved"] == 1
    assert summary["unfilled"] == 1
    assert summary["failed"] == 0


# ---------------------------------------------------------------------------
# template_merge_orchestrator — run_template_merge (fully mocked DB + pipeline steps)
# ---------------------------------------------------------------------------

def make_mock_job():
    return SimpleNamespace(
        id=1, brand_id=1, template_asset_id=10, knowledge_filename="doc.pdf",
        prompt="Make it good", status="pending", current_step=None, progress=0,
        output_path=None, error_detail=None, display_name=None, updated_at=None,
        merge_report=None,
    )


@pytest.mark.unit
def test_run_template_merge_happy_path_sets_completed_and_report():
    job = make_mock_job()
    template_asset = SimpleNamespace(local_path="templates/deck.pptx")
    fake_report = {"slides": [], "summary": {"rewritten": 1}}

    mock_db = MagicMock()
    mock_db.query.return_value.get.side_effect = lambda id_: (
        job if id_ == job.id else template_asset
    )

    with patch.object(orch_mod, "SessionLocal", return_value=mock_db), \
         patch.object(orch_mod.TemplateMergeConfig, "from_db", return_value=TemplateMergeConfig()), \
         patch.object(orch_mod, "resolve_storage", return_value="/tmp/deck.pptx"), \
         patch("os.path.isfile", return_value=True), \
         patch.object(orch_mod, "analyze_template", return_value=["profile1"]), \
         patch.object(orch_mod, "plan_deck", return_value=None), \
         patch.object(orch_mod, "generate_slide_contents", return_value=[{"1": "text"}]), \
         patch.object(orch_mod, "render_merged_pptx", return_value=("/tmp/out.pptx", fake_report)), \
         patch.object(orch_mod, "job_dir", return_value="/tmp/jobs/tm_1"), \
         patch.object(orch_mod, "to_relative", return_value="jobs/tm_1/deck_merged.pptx"):
        orch_mod.run_template_merge(job_id=1)

    assert job.status == "completed"
    assert job.progress == 100
    assert job.output_path == "jobs/tm_1/deck_merged.pptx"
    assert job.merge_report == fake_report


@pytest.mark.unit
def test_run_template_merge_missing_job_logs_and_returns():
    mock_db = MagicMock()
    mock_db.query.return_value.get.return_value = None

    with patch.object(orch_mod, "SessionLocal", return_value=mock_db):
        orch_mod.run_template_merge(job_id=999)  # must not raise

    mock_db.close.assert_called_once()


@pytest.mark.unit
def test_run_template_merge_missing_template_asset_marks_error():
    job = make_mock_job()

    mock_db = MagicMock()
    mock_db.query.return_value.get.side_effect = lambda id_: (
        job if id_ == job.id else None
    )
    mock_db2 = MagicMock()
    mock_db2.query.return_value.get.return_value = job

    with patch.object(orch_mod, "SessionLocal", side_effect=[mock_db, mock_db2]), \
         patch.object(orch_mod.TemplateMergeConfig, "from_db", return_value=TemplateMergeConfig()):
        orch_mod.run_template_merge(job_id=1)

    assert job.status == "error"
    assert "not found" in job.error_detail


@pytest.mark.unit
def test_run_template_merge_missing_file_on_disk_marks_error():
    job = make_mock_job()
    template_asset = SimpleNamespace(local_path="templates/deck.pptx")

    mock_db = MagicMock()
    mock_db.query.return_value.get.side_effect = lambda id_: (
        job if id_ == job.id else template_asset
    )
    mock_db2 = MagicMock()
    mock_db2.query.return_value.get.return_value = job

    with patch.object(orch_mod, "SessionLocal", side_effect=[mock_db, mock_db2]), \
         patch.object(orch_mod.TemplateMergeConfig, "from_db", return_value=TemplateMergeConfig()), \
         patch.object(orch_mod, "resolve_storage", return_value=None):
        orch_mod.run_template_merge(job_id=1)

    assert job.status == "error"
    assert "not found on disk" in job.error_detail


@pytest.mark.unit
def test_run_template_merge_render_failure_marks_error_not_completed():
    job = make_mock_job()
    template_asset = SimpleNamespace(local_path="templates/deck.pptx")

    mock_db = MagicMock()
    mock_db.query.return_value.get.side_effect = lambda id_: (
        job if id_ == job.id else template_asset
    )
    mock_db2 = MagicMock()
    mock_db2.query.return_value.get.return_value = job

    with patch.object(orch_mod, "SessionLocal", side_effect=[mock_db, mock_db2]), \
         patch.object(orch_mod.TemplateMergeConfig, "from_db", return_value=TemplateMergeConfig()), \
         patch.object(orch_mod, "resolve_storage", return_value="/tmp/deck.pptx"), \
         patch("os.path.isfile", return_value=True), \
         patch.object(orch_mod, "analyze_template", return_value=["profile1"]), \
         patch.object(orch_mod, "plan_deck", return_value=None), \
         patch.object(orch_mod, "generate_slide_contents", return_value=[{"1": "text"}]), \
         patch.object(orch_mod, "render_merged_pptx", side_effect=RuntimeError("render exploded")), \
         patch.object(orch_mod, "job_dir", return_value="/tmp/jobs/tm_1"):
        orch_mod.run_template_merge(job_id=1)

    assert job.status == "error"
    assert "render exploded" in job.error_detail


@pytest.mark.unit
def test_set_status_updates_job_fields_and_commits():
    job = make_mock_job()
    mock_db = MagicMock()

    orch_mod._set_status(mock_db, job, "processing", "Doing work...", 42)

    assert job.status == "processing"
    assert job.current_step == "Doing work..."
    assert job.progress == 42
    mock_db.commit.assert_called_once()
