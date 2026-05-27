import os
import sys
import time
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Add backend directory to sys.path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from database import SessionLocal
import models
from services.premium_layout_engine import classify_slide_shapes, smart_replace_text

def test_in_place_mapping():
    db = SessionLocal()
    try:
        # Load a completed job to get its slides
        # We will use Job 80 which has 20 slides synthesized in status "planned"
        job_id = 80
        job = db.query(models.GenerationJob).get(job_id)
        if not job:
            print("Job 80 not found.")
            return
            
        slides_db = db.query(models.PresentationSlide).filter(
            models.PresentationSlide.job_id == job_id
        ).order_by(models.PresentationSlide.slide_number.asc()).all()
        
        print(f"Loaded {len(slides_db)} slides from DB.")
        
        template_path = "uploads/Tesco Style_2.pptx"
        if not os.path.exists(template_path):
            print("Template not found.")
            return
            
        prs = Presentation(template_path)
        total_template_slides = len(prs.slides)
        print(f"Loaded template with {total_template_slides} slides.")
        
        # In-place mapping:
        # We will loop over the synthesized slides and map them to template slides.
        # Slide 1 (Cover) -> Template Slide 0
        # Slide N (Closing) -> Template Slide 22 (the very last slide)
        # Slides 2..N-1 (Content) -> Template Slides 1..N-2 sequentially (or skipping dividers if preferred)
        
        mapped_indices = []
        for i, s_db in enumerate(slides_db):
            num = i + 1
            title = s_db.title
            content = s_db.content_json or {}
            bullets = content.get("bullets", [])
            
            # Map index
            if num == 1:
                idx = 0
            elif i == len(slides_db) - 1:
                idx = total_template_slides - 1
            else:
                # Map intermediate slides sequentially (ensure we don't overflow template slides list)
                idx = min(i, total_template_slides - 2)
                
            print(f"Mapping Synthesized Slide {num} -> Template Slide {idx}")
            slide = prs.slides[idx]
            mapped_indices.append(idx)
            
            # Classify shapes
            shapes = classify_slide_shapes(slide)
            
            # Replace Title
            if shapes["title"]:
                smart_replace_text(shapes["title"], title, max_font_size=44)
                print(f"  - Replaced Title: '{title}'")
                
            # Replace Subtitle
            if shapes["subtitle"]:
                sub = content.get("subtitle") or content.get("section_label") or "Strategic Outlook"
                smart_replace_text(shapes["subtitle"], sub, max_font_size=20)
                
            # Replace Bullets
            if shapes["bullets"]:
                if bullets:
                    orig_has_bullets = any(p.text.strip().startswith(("•", "●", "-", "▪", "○")) for p in shapes["bullets"].text_frame.paragraphs)
                    if orig_has_bullets:
                        bullet_str = "\n".join(bullets)
                    else:
                        bullet_str = "\n".join([f"• {b}" for b in bullets])
                    smart_replace_text(shapes["bullets"], bullet_str, max_font_size=18)
                    print(f"  - Replaced Bullets: {len(bullets)} items")
                else:
                    shapes["bullets"].text_frame.text = ""
                    
            # Replace / clean metadata shapes
            for meta_shape in shapes["metadata"]:
                meta_text = meta_shape.text_frame.text.strip()
                new_meta = meta_text
                if "prepared for" in meta_text.lower() or "sales impact" in meta_text.lower():
                    new_meta = f"Strategic Portfolio | Date: {time.strftime('%d.%m.%y')} | Prepared for Tesco"
                elif "executive lead" in meta_text.lower():
                    new_meta = "Executive Lead: PowerAI Strategic Synthesis"
                elif "contact:" in meta_text.lower() or "@l-founders.com" in meta_text.lower():
                    new_meta = "Contact: support@powerai.com"
                elif "thank you" in meta_text.lower() or "your contacts" in meta_text.lower():
                    new_meta = "Thank you. | Strategic Alliance & Next Steps"
                smart_replace_text(meta_shape, new_meta, max_font_size=12)
                
            # Clean Coles / Loblaws / Conad text inside other text frames on the slide
            for sh in slide.shapes:
                if sh.has_text_frame:
                    text = sh.text_frame.text
                    if "coles" in text.lower() or "loblaw" in text.lower() or "conad" in text.lower() or "scrittori" in text.lower():
                        # Replace occurrences with Tesco
                        new_text = text.replace("Coles", "Tesco").replace("coles", "Tesco")
                        new_text = new_text.replace("Loblaws", "Tesco").replace("loblaw", "Tesco")
                        new_text = new_text.replace("Conad", "Tesco").replace("conad", "Tesco")
                        new_text = new_text.replace("Scrittori di Classe", "Clubcard Evolution")
                        
                        # Retain formatting by doing careful paragraph replacements
                        for p in sh.text_frame.paragraphs:
                            p_text = p.text
                            p_text_new = p_text.replace("Coles", "Tesco").replace("coles", "Tesco")
                            p_text_new = p_text_new.replace("Loblaws", "Tesco").replace("loblaw", "Tesco")
                            p_text_new = p_text_new.replace("Conad", "Tesco").replace("conad", "Tesco")
                            p_text_new = p_text_new.replace("Scrittori di Classe", "Clubcard Evolution")
                            p.text = p_text_new
                            
        # Delete template slides that were not mapped
        mapped_set = set(mapped_indices)
        sldIdLst = prs.slides._sldIdLst
        for idx in range(total_template_slides - 1, -1, -1):
            if idx not in mapped_set:
                print(f"Deleting unused Template Slide {idx}")
                del sldIdLst[idx]
                
        output_filename = "Portfolio_Premium_InPlace_Test.pptx"
        output_path = os.path.join("uploads", output_filename)
        prs.save(output_path)
        print(f"SUCCESS! Presentation saved to: {output_path}")
        print(f"File size: {os.path.getsize(output_path) / (1024*1024):.2f} MB")
        
    finally:
        db.close()

if __name__ == "__main__":
    test_in_place_mapping()
