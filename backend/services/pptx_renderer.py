import os
import re
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# --- UNIVERSAL GEOMETRY HELPERS ---
def hex_to_rgb(hex_str: str) -> RGBColor:
    if not isinstance(hex_str, str): return RGBColor(30, 30, 30)
    h = hex_str.lstrip('#')
    if len(h) == 3: h = ''.join([c*2 for c in h])
    if len(h) != 6: return RGBColor(30, 30, 30)
    try: return RGBColor(*tuple(int(h[i:i+2], 16) for i in (0, 2, 4)))
    except ValueError: return RGBColor(30, 30, 30)

def clean_text(text: str) -> str:
    if not isinstance(text, str): return str(text)
    text = re.sub(r'[*_]{1,3}', '', text)
    return text.strip()

# ──────────────────────────────────────────────
# HELPERS PARA ESTILOS AVANZADOS
# ──────────────────────────────────────────────

def set_shape_transparency(shape, opacity: float):
    """Ajusta la transparencia de un shape (0.0 a 1.0)."""
    if opacity >= 1.0: return
    
    # El valor de alpha en OOXML va de 0 a 100000
    alpha_val = int(opacity * 100000)
    
    # Acceder al XML del elemento de relleno
    sp_tree = shape._element
    solid_fill = sp_tree.find('.//' + qn('a:solidFill'))
    if solid_fill is not None:
        srgb = solid_fill.find(qn('a:srgbClr'))
        if srgb is not None:
            # Eliminar alpha existente si hay
            existing_alpha = srgb.find(qn('a:alpha'))
            if existing_alpha is not None:
                srgb.remove(existing_alpha)
            
            alpha_el = etree.SubElement(srgb, qn('a:alpha'))
            alpha_el.set('val', str(alpha_val))


def apply_corner_style(shape, style: str):
    """Ajusta el radio de las esquinas según el estilo."""
    if style == "rounded":
        # MSO_SHAPE.ROUNDED_RECTANGLE es 5
        # Podemos ajustar el radio si es necesario vía shape.adjustments
        try:
            if hasattr(shape, "adjustments") and len(shape.adjustments) > 0:
                shape.adjustments[0] = 0.05 # 5% de radio
        except: pass
    elif style == "pill":
        try:
            if hasattr(shape, "adjustments") and len(shape.adjustments) > 0:
                shape.adjustments[0] = 0.5 # Totalmente redondeado
        except: pass


# ──────────────────────────────────────────────
# ELEMENT RENDERERS
# ──────────────────────────────────────────────


def render_background_image(slide, element, asset_map, sw_in, sh_in):
    """
    Renderiza la imagen de fondo con estrategia 'Cover' usando CROP interno (v12.0).
    Garantiza 0 desbordamiento fuera del canvas.
    """
    img_basename = element.get("source")
    raw_path = asset_map.get(img_basename)
    
    img_path = None
    if raw_path and os.path.exists(raw_path):
        img_path = raw_path
    
    if not img_path:
        # Fallback a directorio uploads
        potential = os.path.join("uploads", str(img_basename))
        if os.path.exists(potential):
            img_path = potential

    if img_path and os.path.exists(img_path):
        try:
            # 1. Obtener dimensiones para calcular el CROP
            with PILImage.open(img_path) as img:
                iw, ih = img.size
            
            sw, sh = Inches(sw_in), Inches(sh_in)
            
            # 2. Añadir imagen ajustada al slide (0,0)
            pic = slide.shapes.add_picture(img_path, 0, 0, sw, sh)
            
            # 3. Calcular y aplicar CROP semántico (Cover strategy)
            # Queremos que la imagen llene el slide manteniendo su aspect ratio
            aspect_slide = sw_in / sh_in
            aspect_img = iw / ih
            
            if aspect_img > aspect_slide:
                # Imagen más ancha: recortar laterales
                total_crop = 1.0 - (aspect_slide / aspect_img)
                pic.crop_left = total_crop / 2
                pic.crop_right = total_crop / 2
            else:
                # Imagen más alta: recortar arriba/abajo
                total_crop = 1.0 - (aspect_img / aspect_slide)
                pic.crop_top = total_crop / 2
                pic.crop_bottom = total_crop / 2
                
            # 4. Enviar al fondo absoluto del árbol XML (Z-Order)
            # Posición 2 suele ser después de los elementos base del layout
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
            
        except Exception as e:
            print(f"  [Renderer] Background Image failed: {e}")
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(sw_in), Inches(sh_in))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(30, 30, 30)
    else:
        print(f"  [Renderer] Warning: Background asset not found for {img_basename}")


def render_shape(slide, element, sx, sy):
    geo = element.get("geometry", {"left": 0, "top": 0, "width": 10, "height": 10})
    style = element.get("style", {})
    
    shape_type_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "ellipse": MSO_SHAPE.OVAL,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    }
    
    s_type = shape_type_map.get(element.get("shape_type", "rectangle"), MSO_SHAPE.RECTANGLE)
    if style.get("corner_style") in ["rounded", "pill"]:
        s_type = MSO_SHAPE.ROUNDED_RECTANGLE

    shape = slide.shapes.add_shape(
        s_type, 
        sx(geo["left"]), sy(geo["top"]), 
        sx(geo["width"]), sy(geo["height"])
    )
    
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(style.get("color", "#CCCCCC"))
    
    # Transparencia (Glassmorphism / Overlays)
    opacity = style.get("opacity", 1.0)
    set_shape_transparency(shape, opacity)
    
    # Quitar bordes por defecto
    shape.line.fill.background()
    
    # Aplicar radio de esquinas
    if s_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        apply_corner_style(shape, style.get("corner_style", "rounded"))


def render_text(slide, element, slide_data, prs, dna, layout_slug):
    content = clean_text(element.get("content", ""))
    if not content: return

    # 1. Extraer Geometría y Estilo que dictó el Art Director
    role = element.get("role", "body")
    geo = element.get("geometry", {"left": 10.0, "top": 10.0, "width": 80.0, "height": 10.0})
    style = element.get("style", {})
    
    # 2. Funciones helper para porcentajes (sx, sy equivalentes)
    sx = prs.slide_width.inches / 100.0
    sy = prs.slide_height.inches / 100.0
    
    # 3. Crear el cuadro de texto
    text_shape = slide.shapes.add_textbox(
        Inches(geo.get("left", 0) * sx), 
        Inches(geo.get("top", 0) * sy), 
        Inches(geo.get("width", 50) * sx), 
        Inches(geo.get("height", 10) * sy)
    )
    tf = text_shape.text_frame
    tf.word_wrap = True
    
    # 4. Manejar múltiples párrafos (Bullets o líneas)
    lines = content.split('\n')
    base_size = style.get("size", 24)
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.text = line
        
        # Alineación
        align_map = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "left": PP_ALIGN.LEFT}
        p.alignment = align_map.get(style.get("align", "left"), PP_ALIGN.LEFT)
        
        # IMPORTANTE: python-pptx requiere aplicar las fuentes a nivel de 'Run'
        if p.runs:
            for run in p.runs:
                run.font.name = style.get("font", "Arial")
                run.font.size = Pt(base_size)
                run.font.color.rgb = hex_to_rgb(style.get("color", "#000000"))
                run.font.bold = style.get("bold", False)


def render_image(slide, element, asset_map, sx, sy):
    # Soporte para ruta directa (path) o nombre de archivo (source) v16.5
    provided_path = element.get("path")
    img_basename = element.get("source")
    geo = element.get("geometry", {"left": 0, "top": 0, "width": 20, "height": 20})
    
    img_path = None
    print(f"  [Renderer] Attempting to render slide asset: path='{provided_path}', source='{img_basename}'")
    
    # 1. Intentar con provided_path
    if provided_path:
        filename = os.path.basename(provided_path)
        potential_uploads = os.path.join("uploads", filename)
        if os.path.exists(potential_uploads):
            img_path = potential_uploads
            print(f"  [Renderer]   - Success: Found by filename in uploads/: {img_path}")
        elif os.path.exists(provided_path):
            img_path = provided_path
            print(f"  [Renderer]   - Success: Found at provided path: {img_path}")
            
    # 2. Intentar con asset_map y source
    if not img_path and img_basename:
        raw_path = asset_map.get(img_basename)
        if raw_path:
            if os.path.exists(raw_path):
                img_path = raw_path
            else:
                potential = os.path.join("uploads", os.path.basename(raw_path))
                if os.path.exists(potential):
                    img_path = potential
                    
        # 3. Intentar nombre de archivo
        if not img_path:
            potential = os.path.join("uploads", str(img_basename))
            if os.path.exists(potential):
                img_path = potential
            elif os.path.exists(str(img_basename)):
                img_path = str(img_basename)

    if not img_path:
        print(f"  [Renderer]   - CRITICAL: No file found for asset {provided_path or img_basename} in any location.")

    if img_path and os.path.exists(img_path):
        try:
            # Aspect Ratio Protection (v60.0)
            with PILImage.open(img_path) as img:
                iw, ih = img.size
            
            target_w_emu = sx(geo["width"]).emu
            target_h_emu = sy(geo["height"]).emu
            
            # Fit strategy
            scale = min(target_w_emu / iw, target_h_emu / ih)
            nw, nh = int(iw * scale), int(ih * scale)
            
            # Center within target box
            ox = sx(geo["left"]).emu + (target_w_emu - nw) // 2
            oy = sy(geo["top"]).emu + (target_h_emu - nh) // 2
            
            slide.shapes.add_picture(img_path, ox, oy, nw, nh)
            print(f"  [Renderer]   - RENDERED SUCCESSFULLY.")
        except Exception as e:
            print(f"  [Renderer]   - Error during render: {e}")


# ──────────────────────────────────────────────
# MAIN RENDERER
# ──────────────────────────────────────────────

def _render_table_v1(slide, element, sx, sy):
    """
    Dibuja una tabla profesional en PPTX a partir de una matriz de datos.
    """
    rows_data = element.get("data", [])
    if not rows_data: return
    
    rows = len(rows_data)
    cols = len(rows_data[0])
    geo = element.get("geometry", {"top": 40, "left": 10, "width": 80, "height": 40})
    
    left = sx(geo["left"])
    top = sy(geo["top"])
    width = sx(geo["width"])
    height = sy(geo["height"])
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # TIPOGRAFÍA DINÁMICA DE CELDA (v23.0)
    # Evitar desbordamiento vertical escaneando la celda más larga
    max_len = max([len(str(c)) for r in rows_data for c in r] + [0])
    base_f_size = 14
    if max_len > 100: base_f_size = 12
    if max_len > 150: base_f_size = 10
    if max_len > 250: base_f_size = 9

    # Estilo de Celda (v16.3 - Corporate Style)
    for r_idx, row in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            
            # Formatear texto de la celda
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(base_f_size)
            para.font.bold = (r_idx == 0) # Header bold
            para.font.name = "Helvetica Neue"
            
            # Color de fondo (Header vs Body)
            fill = cell.fill
            fill.solid()
            if r_idx == 0:
                fill.fore_color.rgb = RGBColor(0, 82, 163) # Tesco Blue
                para.font.color.rgb = RGBColor(255, 255, 255)
            else:
                fill.fore_color.rgb = RGBColor(245, 245, 245)
                para.font.color.rgb = RGBColor(30, 30, 30)

    print(f"  [Renderer]   - TABLE RENDERED ({rows}x{cols}).")


def render_pptx_manifest(design_manifest, asset_map, output_path):
    """
    ENGINE RENDERER v16.1 — ANALYST VISION DRIVEN.
    Aplica el manifest de diseño con soporte para Canvas Ultra-Wide y Decoradores.
    """
    print(f"[Renderer v16.1] Creating presentation: {output_path}...", flush=True)
    
    # ── 1. Inicializar Canvas Dinámico ────────────────────────────────────
    canvas = design_manifest.get("canvas", {})
    slide_w_in = canvas.get("width_inches", 13.33)
    slide_h_in = canvas.get("height_inches", 7.5)
    
    prs = Presentation()
    prs.slide_width  = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)
    blank_layout = prs.slide_layouts[6] 
    
    theme = design_manifest.get("theme", {})
    bg_color_hex = theme.get("background", "#FFFFFF")

    # Helpers de escalado dinámico
    def sx(val): return Inches((val / 100.0) * slide_w_in)
    def sy(val): return Inches((val / 100.0) * slide_h_in)

    for slide_data in design_manifest.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        elements = slide_data.get("elements", [])
        
        # 2. Fondo base
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_color_hex)

        # 3. Z-Order Sorting (v17.2 - Table Priority)
        # Priorizamos tablas y textos al frente, imágenes al fondo.
        def get_layer(el):
            etype = el.get("type")
            erole = el.get("role", "")
            if etype == "image" and erole == "background": return 0
            if etype == "image" and erole == "background_accent": return 1
            if etype == "shape": return 2
            if etype == "image" and erole == "main": return 3
            if etype == "image" and erole == "accent": return 4
            if etype == "image": return 5
            if etype == "table": return 7  # IMPORTANTE: Tablas siempre encima de las imágenes
            if etype == "text": return 8
            if etype == "logo": return 9
            if erole == "person": return 10
            return 1

        if elements is None: elements = []
        elements.sort(key=get_layer)

        for el in elements:
            try:
                el_type = el.get("type")
                role    = el.get("role", "")
                geo     = el.get("geometry", {})
                style   = el.get("style", {})
                
                if el_type == "image" and role == "background":
                    render_background_image_dynamic(slide, el, asset_map, slide_w_in, slide_h_in)
                elif el_type == "shape" or role in ("horizontal_bar", "vertical_bar", "footer_line", "header_zone", "brand_bar", "overlay", "sidebar_zone", "content_panel", "corner_accent", "pill_accent", "brand_dot"):
                    _render_decorator_v2(slide, el, sx, sy)
                elif el_type == "text":
                    _render_text_v2(slide, el, sx, sy)
                elif el_type == "table":
                    _render_table_v1(slide, el, sx, sy)
                elif el_type in ["image", "logo"]:
                    _render_image_v2(slide, el, asset_map, sx, sy)
            except Exception as e:
                print(f"  [Renderer] Error en elemento {el.get('type')}: {e}")

    prs.save(output_path)
    print(f"[Renderer v16.1] PPTX saved successfully.", flush=True)
    return output_path

# ──────────────────────────────────────────────
# DYNAMIC RENDERERS v2 (v16.1)
# ──────────────────────────────────────────────

def render_background_image_dynamic(slide, element, asset_map, sw_in, sh_in):
    img_basename = element.get("source")
    img_path = asset_map.get(img_basename)
    if not (img_path and os.path.exists(img_path)):
        img_path = os.path.join("uploads", str(img_basename))
        
    if img_path and os.path.exists(img_path):
        try:
            with PILImage.open(img_path) as img: iw, ih = img.size
            pic = slide.shapes.add_picture(img_path, 0, 0, Inches(sw_in), Inches(sh_in))
            aspect_slide, aspect_img = sw_in / sh_in, iw / ih
            if aspect_img > aspect_slide:
                tc = 1.0 - (aspect_slide / aspect_img)
                pic.crop_left, pic.crop_right = tc/2, tc/2
            else:
                tc = 1.0 - (aspect_img / aspect_slide)
                pic.crop_top, pic.crop_bottom = tc/2, tc/2
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
        except: pass

def _render_decorator_v2(slide, element, sx, sy):
    geo, style = element.get("geometry", {}), element.get("style", {})
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, sx(geo["left"]), sy(geo["top"]), sx(geo["width"]), sy(geo["height"])
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(style.get("color", "#CCCCCC"))
    set_shape_transparency(shape, style.get("opacity", 1.0))

def _render_text_v2(slide, element, sx, sy):
    content = clean_text(element.get("content", ""))
    if not content: return
    geo, style, role = element.get("geometry", {}), element.get("style", {}), element.get("role", "")
    
    tx = slide.shapes.add_textbox(sx(geo["left"]), sy(geo["top"]), sx(geo["width"]), sy(geo["height"]))
    tf = tx.text_frame
    tf.word_wrap = True
    
    # IMPORTANTE: No usar TEXT_TO_FIT_SHAPE para que respete nuestro font.size
    # tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 
    
    base_size = style.get("size", 18)
    
    for i, line in enumerate(content.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(style.get("align"), PP_ALIGN.LEFT)
        
        # Aplicar estilos a cada RUN para que PowerPoint no los ignore
        if p.runs:
            for run in p.runs:
                run.font.name = style.get("font", "Helvetica Neue")
                run.font.size = Pt(base_size)
                run.font.color.rgb = hex_to_rgb(style.get("color", "#000000"))
                run.font.bold = style.get("bold", role == "title")

def _render_image_v2(slide, element, asset_map, sx, sy):
    geo, role = element.get("geometry", {}), element.get("role", "")
    
    # Soportar ambas nomenclaturas
    provided_path = element.get("path")
    img_basename = element.get("source")
    
    img_path = None
    if provided_path:
        filename = os.path.basename(provided_path)
        if os.path.exists(os.path.join("uploads", filename)):
            img_path = os.path.join("uploads", filename)
        elif os.path.exists(provided_path):
            img_path = provided_path
            
    if not img_path and img_basename:
        raw = asset_map.get(img_basename)
        if raw and os.path.exists(raw): img_path = raw
        elif os.path.exists(os.path.join("uploads", str(img_basename))):
            img_path = os.path.join("uploads", str(img_basename))

    if img_path and os.path.exists(img_path):
        try:
            iw = element.get("width")
            ih = element.get("height")
            if not iw or not ih:
                with PILImage.open(img_path) as img: iw, ih = img.size
            
            tw, th = sx(geo["width"]).emu, sy(geo["height"]).emu
            sc = min(tw / iw, th / ih)
            nw, nh = int(iw * sc), int(ih * sc)
            
            ox = sx(geo["left"]).emu + (tw - nw) // 2
            
            if role == "person":
                oy = sy(geo["top"]).emu + (th - nh)
            else:
                oy = sy(geo["top"]).emu + (th - nh) // 2
                
            slide.shapes.add_picture(img_path, ox, oy, nw, nh)
        except Exception as e:
            print(f"  [Renderer] Failed to insert image {img_path}: {e}")

def render_pptx_from_db(job_id: int, asset_map: dict, output_path: str):
    """
    RENDERIZADOR DB-DRIVEN (v18.5).
    Reconstruye el manifiesto directamente desde la tabla presentation_slides.
    Garantiza que la DB sea la única fuente de verdad.
    """
    from database import SessionLocal
    from models import GenerationJob, PresentationSlide, BrandVisualDna
    
    db = SessionLocal()
    job = db.query(GenerationJob).filter(GenerationJob.id == job_id).first()
    if not job:
        db.close()
        raise ValueError(f"Job {job_id} not found")
        
    dna = db.query(BrandVisualDna).filter(BrandVisualDna.brand_id == job.brand_id).order_by(BrandVisualDna.created_at.desc()).first()
    
    # Reconstrucción del Manifiesto Estratégico
    slides_db = db.query(PresentationSlide).filter(PresentationSlide.job_id == job_id).order_by(PresentationSlide.slide_number.asc()).all()
    
    manifest = {
        "theme": {
            "primary": dna.primary_color if dna else "#0052A3",
            "background": dna.background_color if dna else "#FFFFFF",
            "font_main": dna.primary_font if dna else "Arial"
        },
        "canvas": {
            "width_inches": dna.slide_width_inches if dna else 13.33,
            "height_inches": dna.slide_height_inches if dna else 7.5
        },
        "slides": []
    }
    
    for s in slides_db:
        manifest["slides"].append({
            "slide_number": s.slide_number,
            "title": s.title,
            "elements": s.render_elements,
            "layout": s.layout_slug
        })
        
    db.close()
    return render_pptx_manifest(manifest, asset_map, output_path)
