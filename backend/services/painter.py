import os
import random
import json
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# --- DYNAMIC BREATHING ENGINE (v8.35 - ADAPTIVE SPACING) ---

def hex_to_rgb(hex_str: str) -> RGBColor:
    if not hex_str or not isinstance(hex_str, str): return RGBColor(0, 82, 163)
    h = hex_str.lstrip('#')
    if len(h) < 6: h = "0052A3"
    try:
        return RGBColor(*tuple(int(h[i:i+2], 16) for i in (0, 2, 4)))
    except:
        return RGBColor(0, 82, 163)

def get_luminance(rgb):
    return (0.299 * rgb[0] + 0.587 * rgb[1] + 0.114 * rgb[2]) / 255

def get_contrast_text_color(bg_rgb):
    return RGBColor(255, 255, 255) if get_luminance(bg_rgb) < 0.5 else RGBColor(20, 20, 20)

class GammaPainter:
    def __init__(self, brand_style):
        print(f"  [Painter] --- DYNAMIC BREATHING v10.0 (GRID ARCHITECTURE) ---")
        self.brand = brand_style
        self.primary = hex_to_rgb(brand_style.primary_color)
        self.secondary = hex_to_rgb(brand_style.secondary_color)
        self.bg = hex_to_rgb(brand_style.background_color or "#FFFFFF")
        self.main_font = getattr(brand_style, "font_family", "Arial") or "Arial"
        
        # GRID SYSTEM & MARGINS (The Elite Consultancy Standard)
        self.MARGIN_X = 5.0  # 5% left/right margin
        self.MARGIN_Y = 8.0  # 8% top/bottom margin
        self.GUTTER = 2.0    # 2% space between columns
        self.GRID_COLS = 12  # Standard 12-column layout
        
        self.SAFE_BOTTOM = 100.0 - self.MARGIN_Y
        
        bg_lum = get_luminance(self.bg)
        p_lum = get_luminance(self.primary)
        self.title_color = self.primary if abs(bg_lum - p_lum) > 0.35 else get_contrast_text_color(self.bg)

        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        self.blank_layout = self.prs.slide_layouts[6]

    def w(self, pct): return self.prs.slide_width * (pct / 100.0)
    def h(self, pct): return self.prs.slide_height * (pct / 100.0)

    # --- ARCHITECTURAL GRID SYSTEM ---
    def grid_x(self, col, span=1):
        """Returns the X coordinate (in Emu) for the start of a column (0 to GRID_COLS-1)."""
        available_pct = 100.0 - (self.MARGIN_X * 2)
        col_w_pct = (available_pct - (self.GUTTER * (self.GRID_COLS - 1))) / self.GRID_COLS
        x_pct = self.MARGIN_X + (col * (col_w_pct + self.GUTTER))
        return self.w(x_pct)

    def grid_w(self, span):
        """Returns the width (in Emu) spanning a given number of columns."""
        available_pct = 100.0 - (self.MARGIN_X * 2)
        col_w_pct = (available_pct - (self.GUTTER * (self.GRID_COLS - 1))) / self.GRID_COLS
        w_pct = (col_w_pct * span) + (self.GUTTER * (span - 1))
        return self.w(w_pct)

    def resolve_image(self, asset_file, min_res=300):
        if not asset_file: return None
        candidates = []
        if os.path.isabs(asset_file): candidates.append(asset_file)
        filename = os.path.basename(asset_file)
        candidates.extend([os.path.abspath(os.path.join("uploads", filename)), os.path.abspath(os.path.join("backend", "uploads", filename))])
        for target in candidates:
            if os.path.exists(target): return target
        return None

    def add_rect(self, slide, x, y, w, h, color, transparency=0.0, rounded=False):
        shape = slide.shapes.add_shape(5 if rounded else 1, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        if transparency > 0: shape.fill.transparency = transparency 
        shape.line.fill.background()
        return shape

    def add_accent_line(self, slide, x, y, w, h_pt=4, color=None):
        """v10.0: Draws a crisp accent line (e.g., the red line under titles)."""
        color = color or self.secondary
        shape = slide.shapes.add_shape(1, x, y, w, Pt(h_pt))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_text(self, slide, text, x, y, w, h, size=24, color=None, bold=False, italic=False, align=PP_ALIGN.LEFT, v_align=None, margin=None):
        tx = slide.shapes.add_textbox(x, y, w, h)
        tf = tx.text_frame
        tf.word_wrap = True
        
        # v10.0: Strict margin control for precise grid alignment
        if margin is not None:
            tf.margin_left = margin
            tf.margin_right = margin
            tf.margin_top = margin
            tf.margin_bottom = margin
        else:
            tf.margin_left = Inches(0)
            tf.margin_right = Inches(0)
            tf.margin_top = Inches(0)
            tf.margin_bottom = Inches(0)

        if v_align:
            tf.vertical_anchor = v_align

        p = tf.paragraphs[0]
        p.text = str(text)
        p.alignment = align
        p.font.name = self.main_font
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color or get_contrast_text_color(self.bg)
        return tx

    def add_fitted_image(self, slide, img_path, x, y, max_w, max_h, transparent=False):
        try:
            with Image.open(img_path) as img:
                orig_w, orig_h = img.size
                # v16.9: Intelligent Scaling (Fixed Logo Visibility)
                is_logo = "logo" in img_path
                is_design_element = "design_element" in img_path
                
                if is_logo:
                    # Logos should fit their box perfectly
                    ratio = min(max_w / orig_w, max_h / orig_h)
                elif is_design_element:
                    # Design elements stay as accents
                    max_upscale = 1.5
                    ratio = min(max_w / orig_w, max_h / orig_h, max_upscale)
                    ratio *= 0.6
                else:
                    ratio = min(max_w / orig_w, max_h / orig_h)
                
                new_w = orig_w * ratio
                new_h = orig_h * ratio
                # Centering within the max box
                off_x = (max_w - new_w) / 2
                off_y = (max_h - new_h) / 2
                pic = slide.shapes.add_picture(img_path, x + off_x, y + off_y, width=new_w, height=new_h)
                return pic
        except: return None

    def add_typo_composition(self, slide, full_text, sub_char, asset_path, x, y, size=80, color=None):
        parts = full_text.split(sub_char)
        current_x = x
        color = color or self.title_color
        def get_width(t, s): return Pt(s * 0.55 * len(t)) 

        if parts[0]:
            self.add_text(slide, parts[0], current_x, y, get_width(parts[0], size), Pt(size * 1.2), size=size, bold=True, color=color)
            current_x += get_width(parts[0], size)
        
        asset_w = Pt(size * 0.9)
        self.add_fitted_image(slide, asset_path, current_x, y + Pt(size * 0.1), asset_w, asset_w)
        current_x += asset_w
        
        if len(parts) > 1 and parts[1]:
            self.add_text(slide, parts[1], current_x, y, get_width(parts[1], size), Pt(size * 1.2), size=size, bold=True, color=color)

    def paint_custom_canvas(self, slide_data):
        slide = self.secure_slide()
        self.add_rect(slide, 0, 0, self.prs.slide_width, self.prs.slide_height, self.bg)
        elements = slide_data.get("elements", [])
        for el in elements:
            type, x, y, w, h = el.get("type"), self.w(el.get("x", 0)), self.h(el.get("y", 0)), self.w(el.get("w", 0)), self.h(el.get("h", 0))
            if type == "text":
                self.add_text(slide, el.get("content"), x, y, w, h, size=el.get("size", 24), bold=el.get("bold", False), color=hex_to_rgb(el.get("color")), align=PP_ALIGN.CENTER if el.get("align") == "center" else PP_ALIGN.LEFT)
            elif type == "image":
                img = self.resolve_image(el.get("path"))
                if img: self.add_fitted_image(slide, img, x, y, w, h)
            elif type == "typo_substitution":
                img = self.resolve_image(el.get("path"))
                if img: self.add_typo_composition(slide, el.get("text"), el.get("char"), img, x, y, size=el.get("size", 80))
        
    def apply_branding(self, slide, slide_data, logo_path, agency):
        """v10.0: Centralized branding logic (Tesco Logo + L-Founders Signature)."""
        is_title = (slide_data.get("slide_number") == 1)
        if logo_path:
            # v16.9: Reliable Logo Placement (Top Right)
            self.add_fitted_image(slide, logo_path, self.w(82), self.h(3), self.w(15), self.h(8))
        self.add_agency_signature(slide, agency, is_title=is_title)

    def paint_split(self, slide_data):
        slide = self.secure_slide()
        img = self.resolve_image(slide_data.get("primary_asset_path"), 300)
        
        # Grid Setup
        start_col, span_cols = 6, 6
        tx, tw = self.grid_x(start_col), self.grid_w(span_cols)
        
        if img:
            self.add_fitted_image(slide, img, 0, 0, self.w(50), self.prs.slide_height)
        else:
            # v16.9: Robust Visual Fallback for split slides
            fb = blend_colors(self.secondary, self.bg, 0.15)
            self.add_rect(slide, 0, 0, self.w(50), self.prs.slide_height, fb)
            self.add_text(slide, "Strategic Vision", self.w(5), self.h(45), self.w(40), self.h(10), size=32, bold=True, color=self.primary, align=PP_ALIGN.CENTER)
            
        self.add_rect(slide, self.w(50), 0, self.w(50), self.prs.slide_height, self.bg)
        
        title = slide_data.get("title", "")
        approx_lines = (len(title) // 25) + 1
        t_size = 28 if approx_lines == 1 else 22
        title_y, title_h = self.MARGIN_Y + 5.0, approx_lines * 6.0
        self.add_text(slide, title, tx, self.h(title_y), tw, self.h(title_h), size=t_size, bold=True, color=self.title_color)
        
        content_top = title_y + title_h + 4.0
        is_last = slide_data.get("is_last", False)
        is_first = slide_data.get("slide_number") == 1
        
        if is_last:
            agency_data = slide_data.get("agency_branding", {})
            self.add_rect(slide, self.w(50), 0, self.w(50), self.prs.slide_height, self.primary)
            closing_title = slide_data.get("title", "STRATEGIC PARTNERSHIP").upper()
            self.add_text(slide, closing_title, self.w(55), self.h(35), self.w(40), self.h(10), size=28, bold=True, color=RGBColor(255, 255, 255))
            self.add_text(slide, "Strategic Alliance & Next Steps", self.w(55), self.h(46), self.w(40), self.h(6), size=14, color=RGBColor(255, 255, 255), italic=True)
            self.add_rect(slide, self.w(55), self.h(58), self.w(10), Pt(1), self.secondary)
            self.add_text(slide, f"Executive Lead: {agency_data.get('name', 'Strategy Team')}", self.w(55), self.h(62), self.w(40), self.h(5), size=11, bold=True, color=RGBColor(255, 255, 255))
            self.add_text(slide, f"Contact: {agency_data.get('email', 'partners@l-founders.com')}", self.w(55), self.h(67), self.w(40), self.h(5), size=10, color=RGBColor(255, 255, 255))
            self.apply_branding(slide, slide_data, slide_data.get("logo_path"), agency_data)
            return slide

        if is_first:
            meta = slide_data.get("metadata", {})
            prep, date, y_meta = meta.get("prepared_for"), meta.get("date"), content_top + 5.0
            if prep:
                self.add_text(slide, f"Prepared for:", tx, self.h(y_meta), tw, self.h(4), size=10, color=self.secondary, bold=True)
                self.add_text(slide, str(prep), tx, self.h(y_meta + 4), tw, self.h(6), size=14, color=self.title_color); y_meta += 12
            if date:
                self.add_text(slide, f"Date:", tx, self.h(y_meta), tw, self.h(4), size=10, color=self.secondary, bold=True)
                self.add_text(slide, str(date), tx, self.h(y_meta + 4), tw, self.h(6), size=12, color=self.title_color); y_meta += 10
            if meta.get("confidential"):
                self.add_text(slide, "CONFIDENTIAL", tx, self.h(y_meta + 2), tw, self.h(4), size=9, color=self.secondary, bold=True)
            self.apply_branding(slide, slide_data, slide_data.get("logo_path"), slide_data.get("agency_branding", {}))
            return slide

        bullets, metrics = slide_data.get("bullets", []), slide_data.get("metrics", [])
        if metrics and not bullets:
            card_cols = 3
            for idx, m in enumerate(metrics[:4]):
                r, c = idx // 2, idx % 2
                x_pos, y_pos, w_pos = self.grid_x(start_col + (c * card_cols)), self.h(content_top + (r * 22)), self.grid_w(card_cols - 0.5)
                self.add_rect(slide, x_pos, y_pos, w_pos, self.h(20), self.primary, transparency=0.92, rounded=True)
                self.add_text(slide, m.get("value", ""), x_pos, y_pos + self.h(2), w_pos, self.h(10), size=32, bold=True, color=get_contrast_text_color(self.primary), align=PP_ALIGN.CENTER)
                self.add_text(slide, m.get("label", ""), x_pos, y_pos + self.h(12), w_pos, self.h(6), size=11, color=get_contrast_text_color(self.bg), align=PP_ALIGN.CENTER)
        if not bullets and not metrics:
            bullets = ["Strategic focus and incremental value creation.", "Operational excellence through data-driven insights."]
            
        num_b = len(bullets[:5])
        available_h = self.SAFE_BOTTOM - content_top
        row_h = min(8.0, available_h / max(1, num_b))
        for idx, b in enumerate(bullets[:5]):
            y_pct = content_top + (idx * (row_h + 2))
            self.add_rect(slide, tx, self.h(y_pct + 1), Pt(12), Pt(12), self.primary, rounded=True)
            self.add_text(slide, b, tx + Pt(20), self.h(y_pct), tw - Pt(20), self.h(row_h), size=14, color=get_contrast_text_color(self.bg))
        return slide

    def paint_big_metric(self, slide_data):
        metric = str(slide_data.get("metric", "")).strip()
        if not metric:
            metrics = slide_data.get("metrics", [])
            if metrics: metric = metrics[0].get("value", "")
        if len(metric) < 2 and not any(c.isdigit() for c in metric): return self.paint_split(slide_data)
        slide = self.secure_slide()
        self.add_rect(slide, 0, 0, self.prs.slide_width, self.prs.slide_height, self.bg)
        title, tx, tw = slide_data.get("title", ""), self.grid_x(1), self.grid_w(10)
        if title: self.add_text(slide, title, tx, self.h(15), tw, self.h(15), size=36, bold=True, color=self.title_color, align=PP_ALIGN.CENTER)
        self.add_text(slide, metric, 0, self.h(40), self.prs.slide_width, self.h(30), size=140, bold=True, color=self.primary, align=PP_ALIGN.CENTER)
        label = slide_data.get("label", "")
        if label: self.add_text(slide, label, 0, self.h(78), self.prs.slide_width, self.h(10), size=28, align=PP_ALIGN.CENTER, color=get_contrast_text_color(self.bg))
        return slide

    def paint_quote(self, slide_data):
        slide = self.secure_slide()
        self.add_rect(slide, 0, 0, self.prs.slide_width, self.prs.slide_height, self.primary)
        bullets = slide_data.get("bullets") or [""]
        q, tx, tw = bullets[0] if bullets else slide_data.get("title", ""), self.grid_x(2), self.grid_w(8)
        self.add_text(slide, f"\"{q}\"", tx, self.h(30), tw, self.h(50), size=44, bold=True, color=get_contrast_text_color(self.primary), align=PP_ALIGN.CENTER)
        return slide

    def paint_hero(self, slide_data):
        slide = self.secure_slide()
        img = self.resolve_image(slide_data.get("primary_asset_path"), 1024)
        if img:
            self.add_fitted_image(slide, img, 0, 0, self.prs.slide_width, self.prs.slide_height)
            self.add_rect(slide, 0, 0, self.prs.slide_width, self.prs.slide_height, self.primary, transparency=0.7)
        else:
            self.add_rect(slide, 0, 0, self.prs.slide_width, self.prs.slide_height, self.primary)
        tx, tw = self.grid_x(1), self.grid_w(10)
        self.add_text(slide, slide_data.get("title", ""), tx, self.h(35), tw, self.h(20), size=54, bold=True, color=get_contrast_text_color(self.primary), align=PP_ALIGN.CENTER)
        bullets = slide_data.get("bullets", [])
        if bullets: self.add_text(slide, bullets[0], tx, self.h(58), tw, self.h(10), size=24, color=get_contrast_text_color(self.primary), align=PP_ALIGN.CENTER)
        return slide

    def paint_grid(self, slide_data):
        slide = self.secure_slide()
        self.add_rect(slide, 0, 0, self.prs.slide_width, self.prs.slide_height, self.bg)
        tx, tw = self.grid_x(0), self.grid_w(12)
        self.add_text(slide, slide_data.get("title", ""), tx, self.h(6), tw, self.h(15), size=34, bold=True, color=self.title_color)
        bullets, col_span = slide_data.get("bullets", []), 3
        for idx, b in enumerate(bullets[:4]):
            x_pos, w_pos = self.grid_x(idx * col_span), self.grid_w(col_span - 0.5)
            self.add_rect(slide, x_pos, self.h(30), w_pos, self.h(58), self.primary, transparency=0.92, rounded=True)
            self.add_text(slide, b, x_pos + self.w(1), self.h(32), w_pos - self.w(2), self.h(54), size=14, color=get_contrast_text_color(self.bg))
        return slide

    def paint_data_grid_cards(self, slide_data):
        slide = self.secure_slide()
        img = self.resolve_image(slide_data.get("primary_asset_path"), 300)
        if img:
            self.add_fitted_image(slide, img, 0, 0, self.grid_x(5), self.prs.slide_height)
            self.add_rect(slide, self.grid_x(5), 0, self.prs.slide_width - self.grid_x(5), self.prs.slide_height, self.bg)
            start_col, span_cols = 5.5, 6.5
        else:
            self.add_rect(slide, 0, 0, self.prs.slide_width, self.prs.slide_height, self.bg)
            start_col, span_cols = 0, 12
        title = slide_data.get("title", "Strategic Indicators")
        if len(title) > 45: t_size, title_h = 24, 18
        elif len(title) > 25: t_size, title_h = 28, 14
        else: t_size, title_h = 32, 10
        self.add_text(slide, title, self.grid_x(start_col), self.h(self.MARGIN_Y + 3), self.grid_w(span_cols), self.h(title_h), size=t_size, bold=True, color=self.title_color)
        accent_y = self.MARGIN_Y + title_h + 1
        self.add_accent_line(slide, self.grid_x(start_col), self.h(accent_y), self.grid_w(1.5), h_pt=3, color=self.secondary)
        metrics = slide_data.get("metrics", [])
        if not metrics: metrics = [{"label": b, "value": "--"} for b in slide_data.get("bullets", [])]
        num_metrics = len(metrics)
        if num_metrics <= 2: card_cols, card_rows = span_cols // 2, 1
        elif num_metrics <= 4: card_cols, card_rows = span_cols // 2, 2
        else: card_cols, card_rows = span_cols // 3, 2
        y_start, card_h = self.h(self.MARGIN_Y + 16), self.h(22) if card_rows == 2 else self.h(32)
        for idx, m in enumerate(metrics[:6]):
            r, c = idx // (span_cols // card_cols), idx % (span_cols // card_cols)
            x_pos, y_pos, w_pos = self.grid_x(start_col + (c * card_cols)), y_start + (r * (card_h + self.h(self.GUTTER))), self.grid_w(card_cols)
            self.add_rect(slide, x_pos, y_pos, w_pos, card_h, self.primary, rounded=True)
            val, v_len = str(m.get("value", "")), len(str(m.get("value", "")))
            if v_len < 6: v_size = 44
            elif v_len < 12: v_size = 28
            elif v_len < 20: v_size = 18
            else: v_size = 12
            self.add_text(slide, val, x_pos, y_pos + self.h(1), w_pos, self.h(11), size=v_size, bold=True, color=get_contrast_text_color(self.primary), align=PP_ALIGN.CENTER, v_align=MSO_ANCHOR.MIDDLE)
            self.add_text(slide, m.get("label", ""), x_pos + self.w(1), y_pos + self.h(12), w_pos - self.w(2), self.h(8), size=11, color=get_contrast_text_color(self.primary), align=PP_ALIGN.CENTER)
            if m.get("growth"): self.add_text(slide, m.get("growth"), x_pos, y_pos + self.h(18), w_pos, self.h(4), size=10, color=get_contrast_text_color(self.primary), italic=True, align=PP_ALIGN.CENTER)
        bullets = slide_data.get("bullets", [])
        if not img and bullets:
            y_bullets = y_start + (card_rows * (card_h + self.h(self.GUTTER))) + self.h(4)
            for i, b in enumerate(bullets[:3]): self.add_text(slide, b, self.grid_x(start_col), y_bullets + (i * self.h(7)), self.grid_w(span_cols), self.h(6), size=16, color=self.title_color)
        return slide

    def paint_pillars(self, slide_data):
        slide = self.secure_slide()
        self.add_rect(slide, 0, 0, self.prs.slide_width, self.prs.slide_height, self.bg)
        self.add_text(slide, slide_data.get("title", ""), self.grid_x(0), self.h(self.MARGIN_Y + 3), self.grid_w(12), self.h(12), size=38, bold=True, color=self.title_color, align=PP_ALIGN.CENTER)
        self.add_accent_line(slide, self.grid_x(5), self.h(self.MARGIN_Y + 12), self.grid_w(2), h_pt=4, color=self.secondary)
        bullets = slide_data.get("bullets", [])
        num_bullets = min(len(bullets), 4)
        if num_bullets == 0: return slide
        col_span = 10 // num_bullets
        start_col = (12 - (num_bullets * col_span)) // 2
        for idx, b in enumerate(bullets[:num_bullets]):
            x_pos, w_pos = self.grid_x(start_col + (idx * col_span)), self.grid_w(col_span)
            self.add_rect(slide, x_pos + self.w(1), self.h(35), w_pos - self.w(2), self.h(50), self.primary, rounded=True)
            self.add_text(slide, b, x_pos + self.w(2), self.h(38), w_pos - self.w(4), self.h(44), size=16, color=get_contrast_text_color(self.primary), align=PP_ALIGN.CENTER, v_align=MSO_ANCHOR.MIDDLE)
        return slide

    def add_agency_signature(self, slide, agency, is_title=False):
        name, footer_y, footer_w = "founders of loyalty", 92.0, 100.0 - (self.MARGIN_X * 2)
        footer_color = RGBColor(255, 255, 255) if get_luminance(self.bg) < 0.5 else self.primary
        if is_title:
            client_tag = f"FOR {agency.get('client_name', 'INTERNAL').upper()} USE ONLY"
            self.add_text(slide, f"L — {name.upper()}  {client_tag}", self.w(self.MARGIN_X), self.h(footer_y), self.w(60), self.h(4), size=8, color=footer_color, bold=True)
        else:
            self.add_rect(slide, self.w(self.MARGIN_X), self.h(91), self.w(footer_w), Pt(0.2), footer_color, transparency=0.8)
            self.add_text(slide, f"L — {name}  CONFIDENTIAL", self.w(self.MARGIN_X), self.h(92), self.w(40), self.h(3), size=7, color=footer_color)

    def secure_slide(self):
        slide = self.prs.slides.add_slide(self.blank_layout)
        for shape in list(slide.placeholders):
            sp = shape.element
            sp.getparent().remove(sp)
        return slide

    def save(self, path):
        self.prs.save(path)
        return path

    def render_slides(self, content_json):
        slides, logo_path, agency = content_json.get("slides", []), content_json.get("logo_path"), content_json.get("agency_branding", {})
        print(f"  [Painter] Rendering {len(slides)} slides...")
        for i, slide_data in enumerate(slides):
            l_type = slide_data.get("layout_type", "composition_split")
            print(f"    - Slide {i+1}: {l_type}")
            slide = None
            if l_type == "composition_hero": slide = self.paint_hero(slide_data)
            elif l_type == "composition_split": slide = self.paint_split(slide_data)
            elif l_type == "big_metric": slide = self.paint_big_metric(slide_data)
            elif l_type == "composition_quote": slide = self.paint_quote(slide_data)
            elif l_type == "composition_pillars": slide = self.paint_pillars(slide_data)
            elif l_type == "data_grid_cards": slide = self.paint_data_grid_cards(slide_data)
            elif l_type == "custom_canvas": slide = self.paint_custom_canvas(slide_data)
            else: slide = self.paint_split(slide_data)
            if slide: self.apply_branding(slide, slide_data, logo_path, agency)

def blend_colors(c1_rgb, c2_rgb, ratio):
    return RGBColor(int(c1_rgb[0] * ratio + c2_rgb[0] * (1 - ratio)), int(c1_rgb[1] * ratio + c2_rgb[1] * (1 - ratio)), int(c1_rgb[2] * ratio + c2_rgb[2] * (1 - ratio)))
