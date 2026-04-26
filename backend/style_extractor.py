import os
import sys
import json
import hashlib
from collections import Counter
from llm_provider import generate_json
from dotenv import load_dotenv

load_dotenv()

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None

def get_hex_from_rgb(r, g, b):
    return f"#{int(r):02x}{int(g):02x}{int(b):02x}".upper()

def extract_pptx_images_recursive(shape, upload_dir):
    """Fase 1: Extracción bruta de todos los binarios de imagen."""
    assets = []
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            f_hash = hashlib.sha256(img.blob).hexdigest()[:12]
            ext = img.ext
            img_filename = f"raw_{f_hash}.{ext}"
            out_path = os.path.join(upload_dir, img_filename)
            if not os.path.exists(out_path):
                with open(out_path, "wb") as f: f.write(img.blob)
            assets.append(out_path)
        except: pass
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            assets.extend(extract_pptx_images_recursive(s, upload_dir))
    return assets

def extract_pptx_style(file_path, client_name="brand", update_callback=None):
    if not Presentation: return {"error": "python-pptx not found."}
    
    prs = Presentation(file_path)
    upload_dir = os.path.dirname(file_path)
    fonts = []
    colors = []
    raw_assets = []
    
    total_slides = len(prs.slides)
    print(f"  [DNA] Phase 1: Raw Extraction...")
    
    # --- PASO 1: EXTRACCIÓN BRUTA ---
    for i, slide in enumerate(prs.slides):
        if update_callback:
            perc = int(((i + 1) / total_slides) * 50) # 50% de la barra para extracción
            update_callback(f"Extracting Raw Assets: Slide {i+1} of {total_slides}...", perc)

        for shape in slide.shapes:
            raw_assets.extend(extract_pptx_images_recursive(shape, upload_dir))
            
            # Estilos de texto (se mantienen igual)
            if not shape.has_text_frame: continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name: fonts.append(run.font.name)
                    try:
                        if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                            r, g, b = run.font.color.rgb
                            colors.append(get_hex_from_rgb(r, g, b))
                    except: pass
                        
    # --- PASO 2: DEDUPLICACIÓN VISUAL Y FILTRADO ---
    print(f"  [DNA] Phase 2: Deduplication and Quality Filtering...")
    if update_callback: update_callback("Performing visual deduplication...", 60)
    
    unique_assets = []
    seen_fingerprints = set() # (size_kb, dimensions?) -> Para simplificar usaremos (size_kb)
    
    # Deduplicación por Hash (Física) + Tamaño (Visual aproximada)
    # Ordenamos por tamaño para procesar primero las de mayor calidad si hay repetidas
    raw_assets = list(set(raw_assets))
    for a_path in raw_assets:
        if not os.path.exists(a_path): continue
        size = os.path.getsize(a_path)
        
        # Filtro de Calidad: Ignorar basura muy pequeña (< 10KB)
        if size < 10240:
            os.remove(a_path)
            continue
            
        # Huella digital basada en tamaño (si dos archivos tienen el mismo tamaño exacto, 
        # es un duplicado de PowerPoint casi seguro)
        fingerprint = size 
        if fingerprint not in seen_fingerprints:
            seen_fingerprints.add(fingerprint)
            # Renombrar de raw_ a asset_ para marcar como validada
            final_name = os.path.basename(a_path).replace("raw_", "asset_")
            final_path = os.path.join(upload_dir, final_name)
            os.rename(a_path, final_path)
            unique_assets.append(final_path)
        else:
            # Es un duplicado visual/físico, lo borramos
            os.remove(a_path)

    print(f"  [DNA] Extraction complete. {len(unique_assets)} unique assets ready for AI magic.")
    
    font_counts = Counter(fonts)
    color_counts = Counter(c for c in colors if c not in ["#000000", "#FFFFFF"])
    
    return {
        "primary_fonts": [f[0] for f in font_counts.most_common(5)],
        "dominant_colors": [c[0] for c in color_counts.most_common(5)],
        "extracted_assets": unique_assets,
        "type": "two_step_extraction"
    }

def refine_style_with_llm(extracted_data, update_callback=None):
    if update_callback: update_callback("AI Magic: Analyzing unique brand assets...", 90)
    
    prompt = f"""
    Expert Corporate Art Director. 
    Analyze Brand DNA. Output ONLY JSON.
    """
    try:
        llm_json = generate_json(prompt, specialization="design")
        assets = extracted_data.get("extracted_assets", [])
        llm_json["extracted_assets"] = [os.path.basename(p) for p in assets]
        llm_json["raw_extraction"] = extracted_data
        return llm_json
    except:
        return extracted_data

def main(file_path, client_name="brand", update_callback=None):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pptx":
        result = extract_pptx_style(file_path, client_name, update_callback)
    else:
        result = {"error": "Unsupported type"}
        
    if "error" not in result:
        final_style = refine_style_with_llm(result, update_callback)
        output_path = os.path.join(os.path.dirname(file_path), "style.json")
        with open(output_path, "w") as f: json.dump(final_style, f, indent=2)
        return final_style
    return result

if __name__ == "__main__":
    if len(sys.argv) > 1: main(sys.argv[1])
