"""
artistic_essence_service.py — PowerAI
Artistic Essence Extraction: layouts, gestos de diseñador, composición.
"""
import os
import json
import datetime
from typing import Optional, Callable, List

from llm_provider import generate_vision_json
from dotenv import load_dotenv

load_dotenv()

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

MAX_SLIDES_TO_ANALYZE = 30 

def _pdf_to_images(file_path: str, out_dir: str, max_pages: int = MAX_SLIDES_TO_ANALYZE) -> List[str]:
    if not fitz: raise ImportError("PyMuPDF no disponible.")
    abs_out_dir = os.path.abspath(out_dir)
    os.makedirs(abs_out_dir, exist_ok=True)
    doc = fitz.open(file_path)
    image_paths = []
    for i in range(min(len(doc), max_pages)):
        page = doc[i]
        pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
        fname = os.path.join(abs_out_dir, f"_vision_page_{i}.png")
        pix.save(fname)
        image_paths.append(fname)
    doc.close()
    return image_paths

def _pptx_to_images(file_path: str, out_dir: str, max_slides: int = MAX_SLIDES_TO_ANALYZE) -> List[str]:
    # --- 1. Primary High-Fidelity Rendering via LibreOffice + PyMuPDF ---
    import subprocess
    abs_out_dir = os.path.abspath(out_dir)
    os.makedirs(abs_out_dir, exist_ok=True)
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    temp_pdf_path = os.path.join(abs_out_dir, f"_temp_essence_{base_name}.pdf")
    
    try:
        print(f"  [Essence] Rendering PPTX to PDF via LibreOffice...", flush=True)
        # Convert to pdf in out_dir, naming it the temp_pdf_path
        cmd = [
            "libreoffice", "--headless", "--convert-to", "pdf",
            "--outdir", abs_out_dir, file_path
        ]
        subprocess.run(cmd, check=True, capture_output=True, timeout=60)
        
        # Libreoffice output will be named as original base_name + .pdf
        standard_pdf_path = os.path.join(abs_out_dir, base_name + ".pdf")
        if os.path.exists(standard_pdf_path):
            os.rename(standard_pdf_path, temp_pdf_path)
            
        if os.path.exists(temp_pdf_path):
            print(f"  [Essence] Rendering PDF pages to slide PNGs...", flush=True)
            image_paths = _pdf_to_images(temp_pdf_path, abs_out_dir, max_slides)
            
            # Clean up the temporary PDF file
            try:
                os.remove(temp_pdf_path)
            except: pass
            
            if image_paths:
                print(f"  [Essence] Success: Rendered {len(image_paths)} full slides using LibreOffice+PyMuPDF.", flush=True)
                return image_paths
    except Exception as e:
        print(f"  [Essence] Warning: LibreOffice slide rendering failed: {e}. Falling back to extraction...", flush=True)
        # Clean up temp PDF if left over
        if os.path.exists(temp_pdf_path):
            try: os.remove(temp_pdf_path)
            except: pass

    # --- 2. Fallback Legacy Extraction (Extracts embedded thumbnails or large pictures) ---
    image_paths = []
    if not Presentation: return []
    prs = Presentation(file_path)
    for i in range(min(len(prs.slides), max_slides)):
        slide = prs.slides[i]
        fname = os.path.join(abs_out_dir, f"_vision_slide_{i}.png")
        try:
            # 1. Intentar miniatura oficial
            thumb_rel = None
            for rel in slide.part.rels.values():
                if "thumbnail" in rel.reltype.lower():
                    thumb_rel = rel; break
            if thumb_rel:
                with open(fname, "wb") as f: f.write(thumb_rel.target_part.blob)
                if os.path.getsize(fname) > 10240:
                    image_paths.append(fname); continue
            
            # 2. Fallback: Extraer la imagen más grande del slide como referencia
            best_img = None
            max_size = 0
            for shape in slide.shapes:
                if getattr(shape, "shape_type", None) == 13: # Picture
                    img_data = shape.image.blob
                    if len(img_data) > max_size:
                        max_size = len(img_data)
                        best_img = img_data
            
            if best_img and max_size > 10240:
                with open(fname, "wb") as f: f.write(best_img)
                image_paths.append(fname)
        except Exception as e:
            print(f"  [Essence] Error extracting PPTX visual: {e}", flush=True)
            # Fallback Pillow
            from PIL import Image, ImageDraw
            img = Image.new("RGB", (1280, 720), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)
            draw.text((50, 50), f"Slide {i+1}: {getattr(slide, 'title', 'No Title')}", fill=(0, 0, 0))
            img.save(fname)
            image_paths.append(fname)
    return image_paths

def _cleanup_images(image_paths: List[str]):
    for p in image_paths:
        try:
            if os.path.exists(p) and "_vision_" in p: os.remove(p)
        except: pass

VISION_BRAND_EXTRACTION_PROMPT = """
You are a Senior Brand Identity Designer analyzing a brand template.

Your task: Extract the VISUAL DESIGN DNA with extreme precision.

Return JSON with exactly this structure:
{
  "visual_strategy": "A high-level explanation of the brand layout style and tone",
  "structural_archetypes": {
    "persistent_blocks": [
      {
        "decorator_type": "accent_line | brand_bar | background_shape | sidebar",
        "description": "where it is placed and how it looks"
      }
    ]
  },
  "design_gestures": {
    "corner_style": "rounded | sharp | pill",
    "visual_density": "dense | balanced | minimal"
  },
  "composition_rules": {
    "max_img_ratio": 0.25,
    "typography_style": "how titles and bodies are treated visually",
    "color_application": "how dominant colors are used in shapes vs text"
  }
}
"""

def analyze_with_vision(image_paths: List[str], cb: Optional[Callable] = None) -> dict:
    """
    STRATEGIC MULTIMODAL ANALYSIS (v18.7).
    Use Gemini Vision to extract structural DNA from reference slides.
    """
    if not image_paths:
        return {}

    # Carga dinámica del Prompt desde la DB para diseño (v20.0)
    from database import SessionLocal
    import models
    db = SessionLocal()
    config_record = db.query(models.SystemConfig).filter(models.SystemConfig.key == "prompt_brand_designer_v1").first()
    db.close()
    
    if config_record:
        prompt = config_record.value
    else:
        # Fallback de diseño estricto
        prompt = VISION_BRAND_EXTRACTION_PROMPT
    
    try:
        if cb: cb("Esencia Artística — Consultando al Director de Arte Visual (Vision LLM)...", 60)
        # Solo enviamos una muestra representativa para no saturar el contexto (v18.7)
        sample_indices = [0, 1, 2, 5, 10, -1]
        sample_images = [image_paths[i] for i in sample_indices if i < len(image_paths)]
        
        from llm_provider import generate_vision_json
        result = generate_vision_json(prompt, sample_images)
        print(f"[Vision] AI RAW Response:\n{json.dumps(result, indent=2)}", flush=True)
        return result
    except Exception as e:
        print(f"  [Vision] Analysis error: {e}")
        return {
            "visual_strategy": "Fallback: Standard corporate style",
            "structural_archetypes": {"persistent_blocks": []},
            "slide_archetypes": {},
            "design_gestures": {"corner_style": "sharp"},
            "composition_rules": {}
        }

import time

def extract_artistic_essence(file_path: str, upload_dir: str, cb: Optional[Callable] = None, brand_id: Optional[int] = None) -> dict:
    start_total = time.time()
    fn_lower = file_path.lower()
    print(f"\n[Essence] === EXTRACTION START (Brand: {brand_id}) ===", flush=True)
    
    if cb: cb("Esencia Artística — Analizando...", 5)
    image_paths = []
    try:
        start_conv = time.time()
        if fn_lower.endswith((".png", ".jpg", ".jpeg", ".webp")): 
            image_paths = [os.path.abspath(file_path)]
        elif ".pdf" in fn_lower: 
            image_paths = _pdf_to_images(file_path, upload_dir)
        elif ".pptx" in fn_lower: 
            image_paths = _pptx_to_images(file_path, upload_dir)
        
        duration_conv = time.time() - start_conv
        print(f"[Essence] Phase 1: Conversion completed en {duration_conv:.2f}s. ({len(image_paths)} images)", flush=True)
    except Exception as e: 
        print(f"[Essence] Conversion ERROR: {e}", flush=True)
        return {"error": f"Procesamiento fallido: {e}"}

    # REGISTRO DE REFERENCIAS (v20.0: Desactivado para evitar ruido en biblioteca)
    # Vision images are only used for analysis, no se registran como activos estratégicos.
    print(f"[Essence] Phase 2: Structure Analysis (Without library registration to avoid noise)", flush=True)

    # VISION ANALYSIS
    start_vision = time.time()
    print("[Essence] Phase 3: Starting Multimodal Vision Analysis...", flush=True)
    vision_result = analyze_with_vision(image_paths, cb)
    duration_vision = time.time() - start_vision
    
    if vision_result.get("error"):
        print(f"[Essence] ERROR in Vision LLM after {duration_vision:.2f}s.", flush=True)
    else:
        print(f"[Essence] Phase 3: Vision Analysis completed en {duration_vision:.2f}s.", flush=True)

    duration_total = time.time() - start_total
    print(f"[Essence] === EXTRACTION FINISHED EN {duration_total:.2f}s ===\n", flush=True)
    
    return {
        "structural_archetypes": vision_result.get("structural_archetypes", {}),
        "design_gestures":       vision_result.get("design_gestures", {}),
        "composition_rules":     vision_result.get("composition_rules", {}),
        "slide_archetypes":      vision_result.get("slide_archetypes", {}),
        "visual_strategy":       vision_result.get("visual_strategy", ""),
        "art_direction_note":    vision_result.get("visual_strategy", ""), # Mapeo directo
        "raw_vision_response":   vision_result,
    }
