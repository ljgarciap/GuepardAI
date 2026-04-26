import os
import sys
import json
from llm_provider import get_embeddings_batch
from dotenv import load_dotenv
from pptx import Presentation
import concurrent.futures
from sqlalchemy import create_engine, text

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

load_dotenv()

# Database Configuration (Ensuring it uses +psycopg for SQLAlchemy)
DB_URL = os.getenv("DATABASE_URL", "postgresql+psycopg://root:root@localhost:5432/ai_db")
engine = create_engine(DB_URL, pool_pre_ping=True)

def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text("text") + "\n\n"
    return text

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n\n"
    return text

def chunk_text(text, chunk_size=1500, overlap=250):
    chunks = []
    start = 0
    while start < len(text):
        chunks.append(text[start:start + chunk_size])
        start += chunk_size - overlap
    return chunks

def ingest_document(file_path, client_name="Internal", update_callback=None, brand_id=None, is_public=False, document_type="company_knowledge"):
    """
    STABLE INGESTION WITH BRAND SOVEREIGNTY, VISIBILITY & TAXONOMY (v12.0).
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in [".pdf", ".pptx"]:
         raise ValueError(f"Format {ext} not supported.")
    
    source_filename = os.path.basename(file_path)
    print(f"[Worker] Starting ingestion for {client_name} (Brand ID: {brand_id}). Type: {document_type}", flush=True)
         
    if update_callback:
        update_callback(f"Reading document: {source_filename}", 5)
    
    if ext == ".pdf":
        if not fitz: raise ValueError("PyMuPDF missing.")
        full_text = extract_text_from_pdf(file_path)
    else:
        full_text = extract_text_from_pptx(file_path)
        
    chunks = chunk_text(full_text)
    valid_chunks = [c for c in chunks if len(c.strip()) >= 5] 
    total_chunks = len(valid_chunks)
    
    print(f"[Worker] Extracted {len(full_text)} chars. Chunks: {total_chunks}", flush=True)
    
    if total_chunks == 0:
        if update_callback:
            update_callback("No text found in document.", -1)
        return

    if update_callback:
        update_callback(f"Analyzing {total_chunks} tactical fragments...", 10)
        
    batch_size = 50
    inserted_total = 0
    
    for i in range(0, total_chunks, batch_size):
        batch_slice = slice(i, i + batch_size)
        batch_texts = valid_chunks[batch_slice]
        current_batch_idx = (i // batch_size) + 1
        total_batches = (total_chunks + batch_size - 1) // batch_size
        
        if update_callback:
            perc = 10 + int((i / total_chunks) * 85)
            update_callback(f"Indexing (Batch {current_batch_idx}/{total_batches})...", perc)
        
        try:
            batch_embeddings = get_embeddings_batch(batch_texts)
            if not batch_embeddings or len(batch_embeddings) == 0:
                 print(f"  [Failure] Batch {current_batch_idx} returned ZERO embeddings.", flush=True)
                 continue
                 
            # Explicit cast to VECTOR format for pgvector
            with engine.connect() as conn:
                with conn.begin():
                    for text_fragment, emb in zip(batch_texts, batch_embeddings):
                        if emb is None: continue
                        emb_pg = f"[{','.join(map(str, emb))}]"
                        
                        conn.execute(
                            text("""
                            INSERT INTO corporate_knowledge (content, metadata, embedding, brand_id, is_public, source_filename, document_type)
                            VALUES (:content, :metadata, cast(:embedding as vector), :brand_id, :is_public, :source_filename, :document_type)
                            """),
                            {
                                "content": text_fragment,
                                "metadata": json.dumps({
                                    "source": source_filename, 
                                    "client": client_name,
                                    "brand_id": brand_id,
                                    "is_public": is_public,
                                    "document_type": document_type
                                }),
                                "embedding": emb_pg,
                                "brand_id": brand_id,
                                "is_public": 1 if is_public else 0,
                                "source_filename": source_filename,
                                "document_type": document_type
                            }
                        )
                        inserted_total += 1
            print(f"  [DB] Consistently persisted batch {current_batch_idx} for Brand {brand_id}.", flush=True)
            
        except Exception as e:
            print(f"  [Error] Batch {current_batch_idx} failed: {e}", flush=True)
            if update_callback:
                update_callback(f"Fatal Sync Error: {str(e)}", -1)
            raise e

    if update_callback:
        update_callback(f"Synchronization finalized ({inserted_total} items).", 100)
    print(f"[Worker] Knowledge base synced. Total: {inserted_total} records.", flush=True)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python ingest_knowledge.py <file_path> [client_name]")
        sys.exit(1)
    
    p = sys.argv[1]
    c = sys.argv[2] if len(sys.argv) > 2 else "Internal"
    if os.path.exists(p):
        ingest_document(p, c)
