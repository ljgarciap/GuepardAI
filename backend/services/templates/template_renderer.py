"""
template_renderer.py — In-place PPTX editor for the Template Merge Engine.

Opens a COPY of the template (the original is never modified), replaces the
text of every slot identified by TemplateAnalyzer with the LLM-generated
content, and saves the result to output_path. Slots are resolved through the
SAME traversal the analyzer used (template_traversal.collect_text_targets),
so addressing can never diverge between analysis and render.

Visual structure is fully preserved:
  - All images, backgrounds, shapes without text → untouched
  - Formatting is preserved by copying the <a:rPr> XML element from the first
    non-empty run of each paragraph (falling back to the frame's first
    non-empty run) into each replacement run. This preserves ALL formatting
    attributes — font name, size, bold, italic, color (RGB and theme),
    character spacing, etc.
  - Bulleted frames (≥2 non-empty paragraphs) get one generated line per
    original paragraph, reusing each paragraph's own <a:pPr> (bullet char,
    numbering, indent) — extra lines clone the last paragraph's formatting,
    leftover paragraphs are blanked. Single-paragraph frames keep the v1
    behavior: soft line-breaks (<a:br>) within the paragraph.
  - Slots not present in the content map are never touched.

Empty-content policy (config.empty_rewrite_policy):
  - action="rewrite" + empty value → text is blanked ("blank", default) or
    the original kept ("keep"); reported as `unfilled` / `kept_original`.
  - action="adapt" + empty value → original kept (`kept_original`).
  - A slide whose whole generation failed (content map is None) keeps all
    its original text; its slots are reported as `failed`.

Returns (output_path, merge_report) — the report lists one outcome per slot:
  rewritten | adapted | preserved | unfilled | kept_original | failed
"""
import copy
import logging
import os
import re
import shutil
from typing import Dict, List, Optional, Tuple

from services.templates.template_analyzer import SlideProfile
from services.templates.template_config import TemplateMergeConfig
from services.templates.template_traversal import collect_text_targets

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from lxml import etree
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

OUTCOMES = ("rewritten", "adapted", "preserved", "unfilled", "kept_original", "failed")


def render_merged_pptx(
    template_path: str,
    profiles: List[SlideProfile],
    slide_contents: List[Optional[Dict[str, str]]],
    output_path: str,
    config: TemplateMergeConfig,
) -> Tuple[str, dict]:
    """
    Copy template → inject content → save to output_path.
    Returns (output_path, merge_report) on success. Raises on fatal error.

    `slide_contents` must be parallel to `profiles` (same length).
    Each element is a dict slot_key (str) → text, or None when that slide's
    generation failed entirely.
    """
    if not _PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is required for template rendering.")

    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    # Work on a copy — the template asset is immutable
    shutil.copy2(template_path, output_path)

    prs = Presentation(output_path)

    if len(profiles) != len(slide_contents):
        logger.warning(
            f"[TemplateMergeRenderer] profiles/contents length mismatch "
            f"({len(profiles)} vs {len(slide_contents)}) — proceeding with min."
        )

    report_slides = []
    for profile, content_map in zip(profiles, slide_contents):
        slide = prs.slides[profile.slide_idx]
        slot_entries = _merge_slide(slide, profile, content_map, config)
        report_slides.append({
            "slide": profile.slide_idx,
            "slots": slot_entries,
            "preserved_shapes": profile.preserved_shapes,
        })

    prs.save(output_path)
    report = {"slides": report_slides, "summary": _summarize(report_slides)}
    logger.info(
        f"[TemplateMergeRenderer] Saved merged PPTX to: {output_path} "
        f"(summary: {report['summary']})"
    )
    return output_path, report


# ─── private ──────────────────────────────────────────────────────────────────

def _merge_slide(
    slide,
    profile: SlideProfile,
    content_map: Optional[Dict[str, str]],
    config: TemplateMergeConfig,
) -> List[dict]:
    """Apply content to one slide; return one report entry per slot."""
    targets, _ = collect_text_targets(slide, config.group_max_depth)
    by_key = {t.key: t for t in targets}

    entries = []
    for slot in profile.slots:
        entry = {
            "key": slot.slot_key,
            "name": slot.shape_name,
            "role": slot.role,
            "action": slot.action,
        }
        entry["outcome"] = _apply_slot(slot, by_key, content_map, config, profile.slide_idx)
        entries.append(entry)
    return entries


def _apply_slot(slot, by_key, content_map, config, slide_idx) -> str:
    if slot.action == "preserve":
        return "preserved"

    if content_map is None:
        # Whole-slide generation failure → keep everything as it was
        return "failed"

    target = by_key.get(slot.slot_key)
    if target is None:
        logger.warning(
            f"[TemplateMergeRenderer] slide {slide_idx} slot '{slot.slot_key}' "
            f"('{slot.shape_name}') not resolvable at render time."
        )
        return "failed"

    # Strip markdown defensively (belt-and-suspenders with template_content)
    value = _strip_markdown((content_map.get(slot.slot_key) or "").strip())

    if not value:
        if slot.action == "rewrite" and config.empty_rewrite_policy == "blank":
            try:
                _blank_text_frame(target.text_frame)
                return "unfilled"
            except Exception as exc:
                logger.warning(
                    f"[TemplateMergeRenderer] slide {slide_idx} slot '{slot.slot_key}' "
                    f"blanking failed: {exc}"
                )
                return "failed"
        return "kept_original"

    try:
        _replace_text_frame(target.text_frame, value)
        if config.reset_autofit:
            _reset_autofit(target.text_frame)
        return "rewritten" if slot.action == "rewrite" else "adapted"
    except Exception as exc:
        logger.warning(
            f"[TemplateMergeRenderer] slide {slide_idx} slot '{slot.slot_key}' "
            f"'{slot.shape_name}' text injection failed: {exc}"
        )
        return "failed"


def _summarize(report_slides: List[dict]) -> Dict[str, int]:
    summary = {outcome: 0 for outcome in OUTCOMES}
    for slide_entry in report_slides:
        for slot_entry in slide_entry["slots"]:
            summary[slot_entry["outcome"]] = summary.get(slot_entry["outcome"], 0) + 1
    return summary


def _replace_text_frame(tf, new_text: str) -> None:
    """
    Replace all text in tf with new_text while preserving ALL formatting.

    Two strategies:
      - Frame originally had ≥2 non-empty paragraphs AND the new content has
        ≥2 lines → paragraph-per-line: each line lands in one original
        paragraph, reusing that paragraph's own <a:pPr> and first-run <a:rPr>
        (bullets, numbering, indent survive). Extra lines clone the last
        paragraph; leftover paragraphs are blanked.
      - Otherwise (v1 behavior) → all lines go into the first paragraph
        separated by <a:br> soft breaks, keeping the shape's geometry stable.
    """
    paragraphs = tf.paragraphs
    if not paragraphs:
        return

    lines = [l for l in new_text.split('\n') if l.strip()]
    if not lines:
        lines = [new_text]

    non_empty = [p for p in paragraphs if p.text.strip()]

    if len(non_empty) >= 2 and len(lines) >= 2:
        _fill_paragraph_per_line(paragraphs, non_empty, lines, tf)
        return

    base_rpr = _capture_base_rpr(tf)
    _set_paragraph_text(paragraphs[0], lines, base_rpr)

    # Clear remaining paragraphs — keep them so shape geometry stays stable
    for para in list(paragraphs)[1:]:
        _clear_paragraph(para)


def _fill_paragraph_per_line(paragraphs, non_empty, lines: List[str], tf) -> None:
    """Map generated lines 1:1 onto the frame's non-empty paragraphs."""
    frame_rpr = _capture_base_rpr(tf)
    last_p_elem = None
    last_rpr = frame_rpr

    for i, line_text in enumerate(lines):
        if i < len(non_empty):
            para = non_empty[i]
            rpr = _capture_para_rpr(para)
            if rpr is None:  # lxml elements must not be truth-tested
                rpr = frame_rpr
            p_elem = para._p
            _clear_p_element(p_elem)
            _append_run(p_elem, line_text, rpr)
            last_p_elem, last_rpr = p_elem, rpr
        else:
            # More lines than paragraphs → clone the last one's formatting
            if last_p_elem is None:
                break
            new_p = copy.deepcopy(last_p_elem)
            _clear_p_element(new_p)
            _append_run(new_p, line_text, last_rpr)
            last_p_elem.addnext(new_p)
            last_p_elem = new_p

    # Fewer lines than paragraphs → blank the leftovers (geometry stays stable)
    for para in non_empty[len(lines):]:
        _clear_paragraph(para)


def _blank_text_frame(tf) -> None:
    """Remove all text (every run in every paragraph), keeping the paragraphs."""
    for para in tf.paragraphs:
        _clear_paragraph(para)


def _reset_autofit(tf) -> None:
    """
    Strip stale <a:normAutofit> fontScale/lnSpcReduction (v2 Fase 3).

    PowerPoint stores the computed autofit shrink on the shape; after we swap
    in longer/shorter text that stored scale is stale — text can visibly
    overflow until a manual edit forces a recompute. Removing the attributes
    (keeping the normAutofit element itself) makes PowerPoint recompute on open.
    Tolerant: any XML surprise is logged, never fatal.
    """
    try:
        body_pr = tf._txBody.find(qn('a:bodyPr'))
        if body_pr is None:
            return
        autofit = body_pr.find(qn('a:normAutofit'))
        if autofit is None:
            return
        for attr in ('fontScale', 'lnSpcReduction'):
            if attr in autofit.attrib:
                del autofit.attrib[attr]
    except Exception as exc:
        logger.warning(f"[TemplateMergeRenderer] autofit reset skipped: {exc}")


def _capture_base_rpr(tf):
    """
    Return a standalone copy of the <a:rPr> element from the first non-empty run.
    Returns None if no explicit run properties are found (formatting then inherits
    from the paragraph/shape theme, which is preserved automatically).
    """
    for para in tf.paragraphs:
        rpr = _capture_para_rpr(para)
        if rpr is not None:
            return rpr
    return None


def _capture_para_rpr(para):
    """Standalone copy of the <a:rPr> from the paragraph's first non-empty run."""
    for run in para.runs:
        if run.text.strip():
            rpr_elem = run._r.find(qn('a:rPr'))
            if rpr_elem is not None:
                # etree.fromstring(tostring(...)) creates an orphan deep copy
                return etree.fromstring(etree.tostring(rpr_elem))
    return None


def _set_paragraph_text(para, lines: List[str], base_rpr) -> None:
    """
    Clear existing runs and write `lines` into `para`.

    Each line becomes an <a:r> run; lines are separated by <a:br> (soft return)
    so the vertical spacing stays compact within the original paragraph frame.
    Each run gets a copy of `base_rpr` to preserve formatting.
    """
    p_elem = para._p
    _clear_p_element(p_elem)

    for i, line_text in enumerate(lines):
        if i > 0:
            # Insert a soft line-break before subsequent lines
            br_elem = etree.SubElement(p_elem, qn('a:br'))
            if base_rpr is not None:
                br_elem.append(etree.fromstring(etree.tostring(base_rpr)))

        run = para.add_run()
        run.text = line_text

        if base_rpr is not None:
            r_elem = run._r
            existing_rpr = r_elem.find(qn('a:rPr'))
            if existing_rpr is not None:
                r_elem.remove(existing_rpr)
            r_elem.insert(0, etree.fromstring(etree.tostring(base_rpr)))

    _push_end_para_rpr_last(p_elem)


def _append_run(p_elem, text: str, rpr) -> None:
    """Append an <a:r> with `text` (and a copy of `rpr`) to a raw <a:p> element."""
    r_elem = etree.SubElement(p_elem, qn('a:r'))
    if rpr is not None:
        r_elem.append(etree.fromstring(etree.tostring(rpr)))
    t_elem = etree.SubElement(r_elem, qn('a:t'))
    t_elem.text = text
    _push_end_para_rpr_last(p_elem)


def _push_end_para_rpr_last(p_elem) -> None:
    """<a:endParaRPr> must stay the last child of <a:p> to keep the XML valid."""
    end = p_elem.find(qn('a:endParaRPr'))
    if end is not None:
        p_elem.remove(end)
        p_elem.append(end)


def _clear_paragraph(para) -> None:
    """Remove all runs and soft-breaks from a paragraph, leaving it empty."""
    _clear_p_element(para._p)


def _clear_p_element(p_elem) -> None:
    for elem in p_elem.findall(qn('a:r')):
        p_elem.remove(elem)
    for elem in p_elem.findall(qn('a:br')):
        p_elem.remove(elem)


def _strip_markdown(text: str) -> str:
    """Defensive strip of markdown formatting in case LLM ignored instructions."""
    if not text:
        return text
    text = re.sub(r'\*{1,3}(.+?)\*{1,3}', r'\1', text, flags=re.DOTALL)
    text = re.sub(r'_{1,2}(.+?)_{1,2}', r'\1', text, flags=re.DOTALL)
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'^[\-\*\+]\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'^\d+\.\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'`([^`]+)`', '\\1', text)
    return text.strip()
