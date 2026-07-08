"""
test_template_merge_integration.py — Integration tests for the Template Merge
Engine: the HTTP endpoints under /api/template-merge/ (served by
routers/template_merge.py since v2 Phase 1) against a real test DB
(db_session, rolled back per test) and the full job lifecycle
(PENDING → PROCESSING → COMPLETED|ERROR) running the real orchestrator against
a real in-memory-generated .pptx file (including a group and a table since v2).

LLM calls are mocked (autouse mock_llm_calls fixture + explicit search_rag
mock) — no tokens spent. File I/O is redirected to a temp tree via
storage_tree (same pattern as test_storage_service.py) so no test artifact
ever lands in backend/storage/.

Spec/design backfill: docs/specs/template-merge.md, docs/designs/template-merge.md
"""
import io
import os
from unittest.mock import patch

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches

import models
from services.core import storage_service as st
from services.templates.template_merge_orchestrator import run_template_merge


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture()
def storage_tree(tmp_path, monkeypatch):
    """Redirects all storage_service roots to an isolated temp tree."""
    root = tmp_path / "storage"
    monkeypatch.setattr(st, "STORAGE_ROOT", str(root))
    monkeypatch.setattr(st, "PUBLIC_ROOT", str(root / "public"))
    monkeypatch.setattr(st, "PRIVATE_ROOT", str(root / "private"))
    monkeypatch.setattr(st, "TMP_ROOT", str(root / "tmp"))
    monkeypatch.setattr(st, "LEGACY_UPLOADS", str(tmp_path / "uploads"))
    monkeypatch.setattr(st, "LEGACY_OUTPUTS", str(tmp_path / "outputs"))
    return tmp_path


@pytest.fixture()
def client(db_session, superadmin_headers):
    from main import app, get_db
    app.dependency_overrides[get_db] = lambda: db_session
    try:
        yield TestClient(app, headers=superadmin_headers)
    finally:
        app.dependency_overrides.clear()


def _make_pptx_bytes(title_text="Placeholder title text here", rich=False):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.text_frame.paragraphs[0].add_run().text = title_text
    if rich:
        # v2 structural coverage: a group with a text box and a 1x2 table
        gtb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(4), Inches(1))
        gtb.text_frame.text = (
            "Grouped narrative block, long enough for the classifier to mark it adaptable"
        )
        slide.shapes.add_group_shape([gtb])
        frame = slide.shapes.add_table(1, 2, Inches(1), Inches(4), Inches(5), Inches(1))
        frame.table.cell(0, 0).text = "KPI"
        frame.table.cell(0, 1).text = (
            "Old metric narrative that is long enough to be replaced by generated content"
        )
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _make_template_asset(db, brand_id=None, local_path="templates/test_deck.pptx", category="pptx_template"):
    asset = models.BrandAsset(
        brand_id=brand_id,
        file_hash="testhash",
        local_path=local_path,
        category=category,
        tags=[],
        manual_tags=[],
        description="Test template asset",
        is_public=0,
    )
    db.add(asset)
    db.flush()
    return asset


def _make_job(db, asset_id, brand_id=None, status="pending", **overrides):
    job = models.TemplateMergeJob(
        brand_id=brand_id,
        template_asset_id=asset_id,
        knowledge_filename="doc.pdf",
        prompt="Write about our Q3 results",
        status=status,
        progress=overrides.pop("progress", 0),
        **overrides,
    )
    db.add(job)
    db.flush()
    return job


# ---------------------------------------------------------------------------
# POST /api/template-merge/upload-template
# ---------------------------------------------------------------------------

@pytest.mark.integration
class TestUploadTemplate:

    def test_upload_rejects_non_pptx_file(self, client, storage_tree):
        res = client.post(
            "/api/template-merge/upload-template",
            files={"file": ("deck.pdf", b"not a pptx", "application/pdf")},
        )
        assert res.status_code == 400

    def test_upload_registers_brand_asset(self, client, db_session, sample_brand, storage_tree):
        content = _make_pptx_bytes()
        res = client.post(
            "/api/template-merge/upload-template",
            files={"file": ("deck.pptx", content, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            data={"brand_id": str(sample_brand.id)},
        )
        assert res.status_code == 200
        body = res.json()
        assert body["category"] == "pptx_template"
        assert body["filename"] == "deck.pptx"

        asset = db_session.query(models.BrandAsset).get(body["asset_id"])
        assert asset is not None
        assert asset.category == "pptx_template"
        assert asset.brand_id == sample_brand.id

        physical = st.resolve(asset.local_path, brand_id=sample_brand.id)
        assert physical is not None and os.path.isfile(physical)


# ---------------------------------------------------------------------------
# POST /api/template-merge/jobs
# ---------------------------------------------------------------------------

@pytest.mark.integration
class TestCreateJob:

    def test_create_job_enqueues_task(self, client, db_session, sample_brand):
        asset = _make_template_asset(db_session, brand_id=sample_brand.id)

        with patch("tasks.celery_run_template_merge.delay") as mock_delay:
            res = client.post("/api/template-merge/jobs", json={
                "template_asset_id": asset.id,
                "knowledge_filename": "doc.pdf",
                "prompt": "Write about our Q3 results",
                "brand_id": sample_brand.id,
            })

        assert res.status_code == 200
        body = res.json()
        assert body["status"] == "pending"
        mock_delay.assert_called_once_with(body["job_id"])

        job = db_session.query(models.TemplateMergeJob).get(body["job_id"])
        assert job is not None
        assert job.knowledge_filename == "doc.pdf"

    def test_create_job_404_when_asset_missing(self, client, db_session):
        with patch("tasks.celery_run_template_merge.delay"):
            res = client.post("/api/template-merge/jobs", json={
                "template_asset_id": 999999,
                "knowledge_filename": "doc.pdf",
                "prompt": "prompt",
            })
        assert res.status_code == 404

    def test_create_job_404_when_asset_not_pptx_template_category(self, client, db_session, sample_brand):
        asset = _make_template_asset(db_session, brand_id=sample_brand.id, category="photos")
        with patch("tasks.celery_run_template_merge.delay"):
            res = client.post("/api/template-merge/jobs", json={
                "template_asset_id": asset.id,
                "knowledge_filename": "doc.pdf",
                "prompt": "prompt",
            })
        assert res.status_code == 404


# ---------------------------------------------------------------------------
# GET /api/template-merge/jobs/{job_id}
# ---------------------------------------------------------------------------

@pytest.mark.integration
class TestJobStatus:

    def test_status_404_when_job_missing(self, client):
        res = client.get("/api/template-merge/jobs/999999")
        assert res.status_code == 404

    def test_status_pending_job_has_no_output_url(self, client, db_session, sample_brand):
        asset = _make_template_asset(db_session, brand_id=sample_brand.id)
        job = _make_job(db_session, asset.id, brand_id=sample_brand.id)

        res = client.get(f"/api/template-merge/jobs/{job.id}")
        assert res.status_code == 200
        body = res.json()
        assert body["status"] == "pending"
        assert body["output_url"] is None

    def test_status_completed_job_exposes_output_url(self, client, db_session, sample_brand, storage_tree):
        asset = _make_template_asset(db_session, brand_id=sample_brand.id)
        out_dir = st.job_dir("tm_status_test")
        out_path = os.path.join(out_dir, "result.pptx")
        with open(out_path, "wb") as f:
            f.write(_make_pptx_bytes())

        job = _make_job(
            db_session, asset.id, brand_id=sample_brand.id,
            status="completed", progress=100, output_path=st.to_relative(out_path),
        )

        res = client.get(f"/api/template-merge/jobs/{job.id}")
        assert res.status_code == 200
        body = res.json()
        assert body["status"] == "completed"
        assert body["output_url"] is not None

    def test_status_exposes_merge_report_and_summary(self, client, db_session, sample_brand):
        asset = _make_template_asset(db_session, brand_id=sample_brand.id)
        report = {"slides": [], "summary": {"rewritten": 3, "unfilled": 1}}
        job = _make_job(
            db_session, asset.id, brand_id=sample_brand.id,
            status="completed", progress=100, merge_report=report,
        )

        res = client.get(f"/api/template-merge/jobs/{job.id}")
        assert res.status_code == 200
        body = res.json()
        assert body["merge_report"] == report
        assert body["merge_summary"] == {"rewritten": 3, "unfilled": 1}

    def test_status_merge_report_null_on_pre_v2_jobs(self, client, db_session, sample_brand):
        asset = _make_template_asset(db_session, brand_id=sample_brand.id)
        job = _make_job(db_session, asset.id, brand_id=sample_brand.id)

        res = client.get(f"/api/template-merge/jobs/{job.id}")
        body = res.json()
        assert body["merge_report"] is None
        assert body["merge_summary"] is None


# ---------------------------------------------------------------------------
# GET /api/template-merge/jobs/{job_id}/download
# ---------------------------------------------------------------------------

@pytest.mark.integration
class TestDownloadJob:

    def test_download_404_when_job_missing(self, client):
        res = client.get("/api/template-merge/jobs/999999/download")
        assert res.status_code == 404

    def test_download_409_when_job_not_completed(self, client, db_session, sample_brand):
        asset = _make_template_asset(db_session, brand_id=sample_brand.id)
        job = _make_job(db_session, asset.id, brand_id=sample_brand.id, status="processing")

        res = client.get(f"/api/template-merge/jobs/{job.id}/download")
        assert res.status_code == 409

    def test_download_404_when_output_file_missing_on_disk(self, client, db_session, sample_brand, storage_tree):
        asset = _make_template_asset(db_session, brand_id=sample_brand.id)
        job = _make_job(
            db_session, asset.id, brand_id=sample_brand.id,
            status="completed", progress=100, output_path="jobs/tm_ghost/never_existed.pptx",
        )

        res = client.get(f"/api/template-merge/jobs/{job.id}/download")
        assert res.status_code == 404

    def test_download_returns_file_when_completed(self, client, db_session, sample_brand, storage_tree):
        asset = _make_template_asset(db_session, brand_id=sample_brand.id)
        out_dir = st.job_dir("tm_download_test")
        out_path = os.path.join(out_dir, "final.pptx")
        with open(out_path, "wb") as f:
            f.write(_make_pptx_bytes())

        job = _make_job(
            db_session, asset.id, brand_id=sample_brand.id,
            status="completed", progress=100, output_path=st.to_relative(out_path),
            display_name="My Merged Deck.pptx",
        )

        res = client.get(f"/api/template-merge/jobs/{job.id}/download")
        assert res.status_code == 200
        assert res.content[:2] == b"PK"  # pptx is a zip archive


# ---------------------------------------------------------------------------
# GET /api/template-merge/templates
# ---------------------------------------------------------------------------

@pytest.mark.integration
class TestListTemplates:

    def test_lists_only_pptx_template_category(self, client, db_session, sample_brand):
        _make_template_asset(db_session, brand_id=sample_brand.id, local_path="a.pptx")
        _make_template_asset(db_session, brand_id=sample_brand.id, local_path="photo.png", category="photos")

        res = client.get("/api/template-merge/templates", params={"brand_id": sample_brand.id})
        assert res.status_code == 200
        body = res.json()
        assert len(body) == 1
        assert body[0]["filename"].endswith("a.pptx") or "a.pptx" in body[0]["filename"]

    def test_filters_by_brand_id(self, client, db_session, sample_brand):
        other_brand = models.Brand(name="OtherBrand_Pytest", about="x", core_value="x")
        db_session.add(other_brand)
        db_session.flush()

        _make_template_asset(db_session, brand_id=sample_brand.id, local_path="mine.pptx")
        _make_template_asset(db_session, brand_id=other_brand.id, local_path="theirs.pptx")

        res = client.get("/api/template-merge/templates", params={"brand_id": sample_brand.id})
        body = res.json()
        assert len(body) == 1


# ---------------------------------------------------------------------------
# Full pipeline: run_template_merge against a real DB row + real .pptx file
# ---------------------------------------------------------------------------

@pytest.mark.integration
class TestFullPipeline:

    def test_pipeline_completes_and_produces_downloadable_pptx(self, db_session, sample_brand, storage_tree):
        template_dir = st.brand_assets_dir(sample_brand.id)
        template_path = os.path.join(template_dir, "source_template.pptx")
        with open(template_path, "wb") as f:
            f.write(_make_pptx_bytes(title_text="Q2 Results Overview", rich=True))

        asset = _make_template_asset(
            db_session, brand_id=sample_brand.id, local_path=st.to_relative(template_path),
        )
        job = _make_job(db_session, asset.id, brand_id=sample_brand.id)
        db_session.commit()  # orchestrator opens its own SessionLocal, must see committed data

        # template_content y template_plan importan generate_json/search_rag
        # por NOMBRE, así que el mock global de conftest (que patchea
        # providers.llm_provider) NO los cubre — sin estos patches explícitos
        # el test hace llamadas LLM REALES.
        outline = {
            "language": "en", "tone": "crisp",
            "slides": [{"slide": 1, "topic": "Q2 results", "key_points": ["12% growth"], "rag_query": "Q2 revenue"}],
        }
        with patch("services.templates.template_plan.generate_json", return_value=outline), \
             patch("services.templates.template_plan.search_rag", return_value="Revenue grew 12% in Q2."), \
             patch("services.templates.template_content.generate_json", return_value={}), \
             patch("services.templates.template_content.search_rag", return_value=["Revenue grew 12% in Q2."]), \
             patch("services.templates.template_merge_orchestrator.SessionLocal", return_value=db_session), \
             patch.object(db_session, "close", lambda: None):
            run_template_merge(job.id)

        db_session.refresh(job)
        assert job.status == "completed"
        assert job.progress == 100
        assert job.output_path is not None

        # v2: the merge report is persisted with per-slot outcomes
        assert job.merge_report is not None
        assert set(job.merge_report.keys()) == {"slides", "summary"}
        reported_keys = {s["key"] for sl in job.merge_report["slides"] for s in sl["slots"]}
        assert any(":" in k for k in reported_keys), "no table cell slot reported"
        assert any("/" in k for k in reported_keys), "no group child slot reported"

        physical = st.resolve(job.output_path, brand_id=sample_brand.id)
        assert physical is not None and os.path.isfile(physical)
        with open(physical, "rb") as f:
            assert f.read(2) == b"PK"

    def test_pipeline_marks_error_when_template_file_deleted_after_scheduling(self, db_session, sample_brand, storage_tree):
        asset = _make_template_asset(
            db_session, brand_id=sample_brand.id, local_path="brands/1/assets/never_uploaded.pptx",
        )
        job = _make_job(db_session, asset.id, brand_id=sample_brand.id)
        db_session.commit()

        with patch("services.templates.template_merge_orchestrator.SessionLocal", return_value=db_session), \
             patch.object(db_session, "close", lambda: None):
            run_template_merge(job.id)

        db_session.refresh(job)
        assert job.status == "error"
        assert job.error_detail
