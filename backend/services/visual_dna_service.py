"""
visual_dna_service.py — PowerAI
Extracción de DNA Visual: colores, fuentes, assets físicos.
Herramienta: Programático (fitz / python-pptx) + LLM texto ligero.
v30.0 — Syntax Corrected & Omnivorous Mode.
"""
import os
import json
import uuid
import hashlib
from collections import Counter
from typing import Optional, Callable

from llm_provider import generate_vision_json, generate_json
from dotenv import load_dotenv

load_dotenv()

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


# ──────────────────────────────────────────────
# UTILIDADES
# ──────────────────────────────────────────────

def _hex(r, g, b) -> str:
    return f"#{int(r):02X}{int(g):02X}{int(b):02X}"


def _is_neutral(hex_color: str) -> bool:
    """Descarta blancos, negros y grises muy puros."""
    h = hex_color.upper().lstrip("#")
    if len(h) != 6:
        return True
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    brightness = (r + g + b) / 3
    spread = max(r, g, b) - min(r, g, b)
    return brightness > 230 or brightness < 20 or spread < 15


# ──────────────────────────────────────────────
# EXTRACCIÓN PROGRAMÁTICA — PDF
# ──────────────────────────────────────────────

def extract_pdf_dna(file_path: str, source_filename: str,
                    upload_dir: str, cb: Optional[Callable] = None) -> dict:
    if not fitz:
        return {"error": "PyMuPDF no disponible."}

    doc = fitz.open(file_path)
    fonts, colors, assets = [], [], []
    seen_hashes = set()
    total = len(doc)

    for i, page in enumerate(doc):
        if cb:
            cb(f"DNA Visual — PDF página {i+1}/{total}", int((i + 1) / total * 90))

        # Texto → fuentes y colores
        for block in page.get_text("dict").get("blocks", []):
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    if span.get("font"):
                        fonts.append(span["font"])
                    c = span.get("color", 0)
                    r, g, b = (c >> 16) & 255, (c >> 8) & 255, c & 255
                    hex_c = _hex(r, g, b)
                    if not _is_neutral(hex_c):
                        colors.append(hex_c)

        # Imágenes nativas con Hashing
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            try:
                base_img = doc.extract_image(xref)
                data = base_img.get("image", b"")
                ext = base_img.get("ext", "png")
                if len(data) > 5120:  # > 5KB
                    f_hash = hashlib.sha256(data).hexdigest()[:16]
                    if f_hash in seen_hashes: continue
                    seen_hashes.add(f_hash)
                    
                    fname = f"pdf_img_{f_hash}.{ext}"
                    out = os.path.join(upload_dir, fname)
                    
                    width, height = 0, 0
                    try:
                        import io
                        from PIL import Image
                        with Image.open(io.BytesIO(data)) as img:
                            width, height = img.size
                    except: pass

                    if not os.path.exists(out):
                        with open(out, "wb") as f: f.write(data)
                    assets.append({"path": fname, "width": width, "height": height})
            except Exception as e:
                print(f"  [DNA] Error extrayendo imagen PDF xref {xref}: {e}", flush=True)

    font_freq = Counter(fonts)
    color_freq = Counter(c for c in colors)

    labeled_assets = {"photos": [], "logos": [], "icons": []}
    for asset_name in assets:
        category = "photos"
        if "logo" in asset_name.lower(): category = "logos"
        elif "icon" in asset_name.lower(): category = "icons"
        labeled_assets[category].append({"path": asset_name, "description": "Extracted PDF asset"})

    return {
        "primary_fonts": [f[0] for f in font_freq.most_common(5)],
        "dominant_colors": [c[0] for c in color_freq.most_common(8)],
        "extracted_assets": labeled_assets,
        "source_type": "pdf",
    }


# ──────────────────────────────────────────────
# EXTRACCIÓN PROGRAMÁTICA — PPTX (Omnivorous v30.0)
# ──────────────────────────────────────────────

def extract_pptx_dna(file_path: str, source_filename: str,
                     upload_dir: str, cb: Optional[Callable] = None) -> dict:
    if not Presentation:
        return {"error": "python-pptx no disponible."}

    from pptx.enum.dml import MSO_FILL
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    fonts, colors, raw_assets_paths = [], [], []
    seen_hashes = set()
    total_slides = len(prs.slides)

    def _process_img_data(blob, ext, prefix):
        if not blob or len(blob) < 4096: return None
        f_hash = hashlib.sha256(blob).hexdigest()[:16]
        if f_hash in seen_hashes: return None
        
        seen_hashes.add(f_hash)
        fname = f"{prefix}_{f_hash}.{ext}"
        out_path = os.path.join(upload_dir, fname)
        
        width, height = 0, 0
        try:
            import io
            from PIL import Image
            with Image.open(io.BytesIO(blob)) as img:
                width, height = img.size
        except: pass

        if not os.path.exists(out_path):
            with open(out_path, "wb") as f: f.write(blob)
        return {"path": fname, "width": width, "height": height}

    def _extract_recursive(shape):
        paths = []
        # 1. Imagen Directa
        if getattr(shape, "shape_type", None) == 13: # PICTURE
            try:
                res = _process_img_data(shape.image.blob, shape.image.ext, "img")
                if res: paths.append(res)
            except: pass
        # 2. Relleno de Imagen
        try:
            if hasattr(shape, "fill") and shape.fill.type == 6: # PICTURE FILL
                img = shape.fill.picture.image
                res = _process_img_data(img.blob, img.ext, "fill")
                if res: paths.append(res)
        except: pass
        # 3. Grupos
        if getattr(shape, "shape_type", None) == 6: # GROUP
            for s in shape.shapes:
                paths.extend(_extract_recursive(s))
        
        # Fuentes y Colores
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name: fonts.append(run.font.name)
                    try:
                        col = run.font.color
                        if col and hasattr(col, "rgb") and col.rgb:
                            colors.append(_hex(*col.rgb))
                    except: pass
        return paths

    # --- 1. MASTERS & LAYOUTS ---
    print(f"  [DNA] Scanning Masters and Layouts (Omnivorous Mode)...", flush=True)
    for master in prs.slide_masters:
        try:
            if master.background.fill.type == 6:
                img = master.background.fill.picture.image
                res = _process_img_data(img.blob, img.ext, "master_bg")
                if res: raw_assets_paths.append(res)
        except: pass
        for shape in master.shapes: raw_assets_paths.extend(_extract_recursive(shape))
        
        for layout in master.slide_layouts:
            try:
                if layout.background.fill.type == 6:
                    img = layout.background.fill.picture.image
                    res = _process_img_data(img.blob, img.ext, "layout_bg")
                    if res: raw_assets_paths.append(res)
            except: pass
            for shape in layout.shapes: raw_assets_paths.extend(_extract_recursive(shape))

    # --- 2. SLIDES ---
    print(f"  [DNA] Scanning Slides...", flush=True)
    for i, slide in enumerate(prs.slides):
        if cb: cb(f"DNA Visual — Slide {i+1}/{total_slides}", int((i+1)/total_slides * 90))
        try:
            if slide.background.fill.type == 6:
                img = slide.background.fill.picture.image
                res = _process_img_data(img.blob, img.ext, "bg")
                if res: raw_assets_paths.append(res)
        except: pass
        for shape in slide.shapes: raw_assets_paths.extend(_extract_recursive(shape))

    # --- 3. CATEGORIZACIÓN ---
    labeled_assets = {"photos": [], "logos": [], "icons": []}
    for asset_info in raw_assets_paths:
        fname = asset_info["path"]
        category = "photos"
        if "logo" in fname.lower(): category = "logos"
        elif "icon" in fname.lower(): category = "icons"
        labeled_assets[category].append({
            "path": fname, 
            "width": asset_info["width"], 
            "height": asset_info["height"],
            "description": "Extracted DNA asset"
        })

    return {
        "primary_fonts": [f[0] for f in Counter(fonts).most_common(5)],
        "dominant_colors": [c[0] for c in Counter(colors).most_common(8)],
        "extracted_assets": labeled_assets,
        "source_type": "pptx",
        "slide_width_inches": round(prs.slide_width.inches, 2),
        "slide_height_inches": round(prs.slide_height.inches, 2),
    }


# ──────────────────────────────────────────────
# REFINAMIENTO CON LLM TEXTO
# ──────────────────────────────────────────────

def refine_dna_with_llm(raw_data: dict, cb: Optional[Callable] = None) -> dict:
    if cb:
        cb("DNA Visual — Sintetizando con LLM...", 92)

    prompt = f"""
You are a Senior Brand Identity Expert. 
Based on programmatically extracted data from an official document:

{json.dumps(raw_data, indent=2)}

Your mission: Synthesize the brand's Visual DNA with high accuracy.

STRICT COLOR RULES:
1. primary_color: This MUST be the main institutional color (e.g., Tesco Blue #00539F). Ignore noise colors from photos or charts.
2. secondary_color: A distinct secondary brand color (e.g., Tesco Red #EE1C2E). Must have high contrast with primary.
3. background_color: Usually #FFFFFF or a very light brand tint.
4. text_main_color: Usually #000000 or very dark grey/navy for readability.
5. accent_color: A vibrant accent color if present, or null.

Return ONLY this JSON:
{{
  "primary_color": "#XXXXXX",
  "secondary_color": "#XXXXXX",
  "background_color": "#XXXXXX",
  "text_main_color": "#XXXXXX",
  "accent_color": "#XXXXXX or null",
  "primary_font": "font name",
  "secondary_font": "font name or null"
}}
"""
    try:
        result = generate_json(prompt, specialization="general")
        result["extracted_assets"] = raw_data.get("extracted_assets", [])
        result["raw_extraction"] = raw_data
        result["slide_width_inches"] = raw_data.get("slide_width_inches", 13.33)
        result["slide_height_inches"] = raw_data.get("slide_height_inches", 7.5)
        return result
    except Exception as e:
        print(f"  [DNA] LLM refinement failed: {e}", flush=True)
        fonts = raw_data.get("primary_fonts", ["Arial"])
        colors = raw_data.get("dominant_colors", ["#333333"])
        return {
            "primary_color": colors[0] if len(colors) > 0 else "#333333",
            "secondary_color": colors[1] if len(colors) > 1 else "#666666",
            "background_color": "#FFFFFF",
            "text_main_color": "#111111",
            "accent_color": colors[2] if len(colors) > 2 else None,
            "primary_font": fonts[0] if fonts else "Arial",
            "secondary_font": fonts[1] if len(fonts) > 1 else None,
            "extracted_assets": raw_data.get("extracted_assets", []),
            "raw_extraction": raw_data,
            "slide_width_inches": raw_data.get("slide_width_inches", 13.33),
            "slide_height_inches": raw_data.get("slide_height_inches", 7.5),
            "llm_error": str(e),
        }


# ──────────────────────────────────────────────
# FUNCIÓN PRINCIPAL
# ──────────────────────────────────────────────

def extract_visual_dna(file_path: str, upload_dir: str,
                       cb: Optional[Callable] = None) -> dict:
    source_filename = os.path.basename(file_path)
    ext = os.path.splitext(file_path)[1].lower()
    print(f"[DNA] Starting extraction for {source_filename} (Type: {ext})", flush=True)

    if cb:
        cb("DNA Visual — Iniciando extracción...", 2)

    if ext == ".pdf":
        raw = extract_pdf_dna(file_path, source_filename, upload_dir, cb)
    elif ext == ".pptx":
        raw = extract_pptx_dna(file_path, source_filename, upload_dir, cb)
    else:
        return {"error": f"Formato no soportado: {ext}"}

    if "error" in raw:
        return raw

    return refine_dna_with_llm(raw, cb)
