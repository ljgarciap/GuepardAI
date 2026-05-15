import sys
import json
import os
import urllib.request
import urllib.parse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import concurrent.futures
import random
import time

# --- AGNOSTIC-CORE RENDERING ENGINE (v3.1 - Precision) ---
# Strictly data-driven. No hardcoded design defaults.

def hex_to_rgb(hex_str: str) -> RGBColor:
    """Safely converts hex to RGBColor. Returns neutral grey on failure."""
    if not isinstance(hex_str, str): return RGBColor(128, 128, 128)
    h = hex_str.lstrip('#')
    if len(h) != 6: return RGBColor(128, 128, 128)
    try: return RGBColor(*tuple(int(h[i:i+2], 16) for i in (0, 2, 4)))
    except ValueError: return RGBColor(128, 128, 128)

def clean_text(val) -> str:
    """Sanitizes text by stripping JSON/dict structures if returned by LLM."""
    if isinstance(val, dict):
        # Extract the logical description or title, or just join values
        if "description" in val: return str(val["description"])
        if "title" in val: return str(val["title"])
        return " ".join(str(v) for v in val.values())
    return str(val)

class Theme:
    """Manages brand visual patterns extracted from DNA."""
    def __init__(self, style_json):
        self.primary_color = hex_to_rgb(style_json.get("primary_color"))
        self.secondary_color = hex_to_rgb(style_json.get("secondary_color"))
        self.bg_color = hex_to_rgb(style_json.get("background_color", "#FFFFFF"))
        self.text_main = hex_to_rgb(style_json.get("text_main_color", "#333333"))
        
        self.font_main = style_json.get("primary_font", "Arial")
        self.font_body = style_json.get("secondary_font", "Calibri")
        
        # AGNOSTIC PATTERNS
        patterns = style_json.get("visual_patterns", {})
        self.density = patterns.get("density", "Low")      # High | Medium | Low
        self.gravity = patterns.get("gravity", "Executive") # Massive | Executive
        self.accent = patterns.get("accent", "None")       # Dot | Line | None

SAFE_BACKUP_IMAGES = [
    "https://images.unsplash.com/photo-1497215842964-222b430dc094?w=1600&q=80",
    "https://images.unsplash.com/photo-1460925895917-afdab827c52f?w=1600&q=80",
    "https://images.unsplash.com/photo-1557804506-669a67965ba0?w=1600&q=80",
    "https://images.unsplash.com/photo-1542744173-8e7e53415bb0?w=1600&q=80"
]

def download_asset(prompt: str, slide_num: int) -> str:
    # Use image_prompt (if structured) or prompt (if simple)
    image_prompt = prompt or "corporate strategic metrics"
    
    # Inject random entropy keywords to break Unsplash caching and force pool diversity
    pool_keywords = ["innovation", "strategy", "futuristic", "abstract", "modern", "visionary", "digital"]
    entropy_word = random.choice(pool_keywords)
    
    # Enhanced query with entropy and unique signature
    search_query = f"{image_prompt}, {entropy_word}"
    sig = f"{slide_num}_{int(time.time())}_{random.randint(0, 1000)}"
    url = f"https://source.unsplash.com/1600x900/?{urllib.parse.quote(search_query)}&sig={sig}"
    
    temp_path = os.path.join("uploads", f"asset_{slide_num}_{os.urandom(2).hex()}.jpg")
    try:
        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=12) as response, open(temp_path, 'wb') as out_file:
            out_file.write(response.read())
        return temp_path
    except:
        fb = SAFE_BACKUP_IMAGES[slide_num % len(SAFE_BACKUP_IMAGES)]
        try:
            with urllib.request.urlopen(fb, timeout=10) as r, open(temp_path, 'wb') as f:
                f.write(r.read())
            return temp_path
        except: return None

def create_slide_builder(content_json_path, style_json_path, output_path):
    print("Initializing Agnostic-Core Engine v3.0...")
    with open(content_json_path, 'r', encoding='utf-8') as f: content = json.load(f)
    with open(style_json_path, 'r', encoding='utf-8') as f: style_data = json.load(f)
    
    theme = Theme(style_data)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    slides_data = content.get("slides", [])
    img_paths = {}
    def fetch(idx, p): return idx, download_asset(p, idx)
    with concurrent.futures.ThreadPoolExecutor(max_workers=8) as ex:
        futures = {ex.submit(fetch, s.get("slide_number", i+1), s.get("image_prompt", s.get("title", ""))): s.get("slide_number", i+1) 
                   for i, s in enumerate(slides_data)}
        for f in concurrent.futures.as_completed(futures):
            idx, path = f.result(); img_paths[idx] = path

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        idx = slide_data.get("slide_number", 1)
        layout = slide_data.get("layout_type", "composition_split")
        img_path = img_paths.get(idx)
        
        title_text = str(slide_data.get("title", ""))
        bullets = slide_data.get("bullets", [])

        # AGNOSTIC BACKGROUND LOGIC:
        # High Density brands get Full Primary backgrounds on key slides (Hero, Metrics)
        is_high_impact = (theme.density == "High" and ("hero" in layout or "metric" in layout))
        
        bg_rgb = theme.primary_color if is_high_impact else theme.bg_color
        text_rgb = RGBColor(255, 255, 255) if is_high_impact else theme.primary_color
        body_rgb = RGBColor(255, 255, 255) if is_high_impact else theme.text_main

        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        bg.fill.solid(); bg.fill.fore_color.rgb = bg_rgb; bg.line.fill.background()

        # --- COMPOSITION HERO ---
        if "hero" in layout:
            if img_path and theme.density != "High": # Favor image in low/med density
                slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(10), Inches(7.5))
                ov = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
                ov.fill.solid(); ov.fill.fore_color.rgb = RGBColor(0,0,0); ov.fill.transparency = 0.5; ov.line.fill.background()
            
            # --- ADAPTIVE TYPOGRAPHY (PIXEL-PERFECT) ---
            # Automatically shrinks if title is too long (>25 chars)
            base_pt = 80 if theme.gravity == "Massive" else 60
            title_size = Pt(max(32, base_pt - (max(0, len(title_text) - 20) * 1.5)))
            
            tx = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(4))
            tf = tx.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = 1; p.text = title_text
            p.font.name = theme.font_main; p.font.size = title_size; p.font.color.rgb = RGBColor(255,255,255); p.font.bold = True
            
            # DOT ACCENT (Agnostic implementation of the Tesco 'Dot')
            if theme.accent == "Dot":
                p.text = p.text.rstrip(".")
                dot = slide.shapes.add_shape(9, Inches(8.5), Inches(3.5), Inches(0.2), Inches(0.2)) # Circle
                dot.fill.solid(); dot.fill.fore_color.rgb = theme.secondary_color; dot.line.fill.background()

        # --- COMPOSITION SPLIT (ALTERNATING SYMMETRY) ---
        elif "split" in layout:
            is_inverted = (idx % 2 == 0) # Alternate left/right based on slide index
            img_x = Inches(0) if is_inverted else Inches(6)
            txt_x = Inches(4.5) if is_inverted else Inches(0.5)
            
            if img_path: slide.shapes.add_picture(img_path, img_x, Inches(0), Inches(4), Inches(7.5))
            
            # Accent Line
            if theme.accent == "Line":
                lx = Inches(4) if is_inverted else Inches(5.8)
                accent = slide.shapes.add_shape(1, lx, Inches(0), Inches(0.2), Inches(7.5))
                accent.fill.solid(); accent.fill.fore_color.rgb = theme.primary_color; accent.line.fill.background()
            
            # Adaptive size for Split titles
            title_size = Pt(max(24, 44 - (max(0, len(title_text) - 18) * 1.2)))
            tx = slide.shapes.add_textbox(txt_x, Inches(0.8), Inches(5.2), Inches(2.5))
            tf = tx.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = title_text
            p.font.name = theme.font_main; p.font.size = title_size; p.font.color.rgb = text_rgb; p.font.bold = True
            
            bt = slide.shapes.add_textbox(txt_x, Inches(3), Inches(5.2), Inches(4))
            btf = bt.text_frame; btf.word_wrap = True
            for b in bullets:
                bp = btf.add_paragraph(); bp.text = "• " + clean_text(b)
                bp.font.name = theme.font_body; bp.font.size = Pt(18); bp.font.color.rgb = body_rgb

        # --- COMPOSITION PILLARS ---
        elif any(x in layout for x in ["pillars", "columns", "cards"]):
            hd = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            tf = hd.text_frame; p = tf.paragraphs[0]; p.text = title_text.upper()
            p.font.name = theme.font_main; p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = text_rgb
            
            items = bullets[:3]
            for i, item in enumerate(items):
                x = Inches(0.5 + (i * 3.1))
                # If high density, use fully colored blocks
                card = slide.shapes.add_shape(1, x, Inches(2), Inches(2.8), Inches(4.5))
                if theme.density == "High":
                    card.fill.solid(); card.fill.fore_color.rgb = theme.secondary_color
                    txt_color = RGBColor(255,255,255)
                else:
                    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(255,255,255)
                    card.line.color.rgb = theme.primary_color; card.line.width = Pt(1.5)
                    txt_color = theme.text_main
                
                txt = slide.shapes.add_textbox(x+Inches(0.1), Inches(2.1), Inches(2.6), Inches(4.2))
                tf = txt.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = clean_text(item)
                p.font.name = theme.font_body; p.font.size = Pt(16); p.font.color.rgb = txt_color

        # --- COMPOSITION IMPACT (Big Metric) ---
        elif layout == "big_metric":
            metric = clean_text(slide_data.get("metric", "100%"))
            label = clean_text(slide_data.get("label", title_text))
            
            # ADAPTIVE FONT SIZE (v23.0)
            # If the metric is a long string (e.g. £1,750 M) we scale down from 160 to fit.
            base_size = 160
            if len(metric) > 5: base_size = 110
            if len(metric) > 10: base_size = 80
            if len(metric) > 15: base_size = 60
            
            mx = slide.shapes.add_textbox(Inches(0), Inches(2), Inches(10), Inches(3))
            mf = mx.text_frame; p = mf.paragraphs[0]; p.alignment = 1; p.text = metric
            p.font.name = theme.font_main; p.font.size = Pt(base_size); p.font.bold = True; p.font.color.rgb = text_rgb
            
            lx = slide.shapes.add_textbox(Inches(0), Inches(5.5), Inches(10), Inches(1))
            lf = lx.text_frame; p = lf.paragraphs[0]; p.alignment = 1; p.text = label.upper()
            p.font.name = theme.font_body; p.font.size = Pt(28); p.font.color.rgb = body_rgb
            
        # --- COMPOSITION QUOTE (Executive Impact) ---
        elif "quote" in layout:
            # Centered text with decorative primary color quotes
            base_pt = 36 if len(title_text) < 50 else 28
            tx = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(3))
            tf = tx.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = 1
            p.text = f'"{title_text}"'
            p.font.name = theme.font_main; p.font.size = Pt(base_pt); p.font.italic = True; p.font.color.rgb = theme.primary_color
            
        # --- COMPOSITION GRID (2x2 Matrix) ---
        elif "grid" in layout:
            hd = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            p = hd.text_frame.paragraphs[0]; p.text = title_text.upper(); p.font.bold = True
            p.font.name = theme.font_main; p.font.size = Pt(28); p.font.color.rgb = text_rgb
            
            items = bullets[:4]
            for i, item in enumerate(items):
                col, row = i % 2, i // 2
                x, y = Inches(0.8 + (col * 4.5)), Inches(1.8 + (row * 2.5))
            # Grid cell card
                card = slide.shapes.add_shape(1, x, y, Inches(4), Inches(2.2))
                card.fill.solid(); card.fill.fore_color.rgb = theme.bg_color
                card.line.color.rgb = theme.secondary_color; card.line.width = Pt(1)
                
                txt = slide.shapes.add_textbox(x+Inches(0.2), y+Inches(0.2), Inches(3.6), Inches(1.8))
                tf = txt.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = clean_text(item)
                p.font.name = theme.font_body; p.font.size = Pt(14); p.font.color.rgb = body_rgb
            
        # --- SAFETY FALLBACK (ANTI-BLANK-SLIDE) ---
        else:
            # If no layout matched, force a default Split composition to ensure text visibility
            if img_path: slide.shapes.add_picture(img_path, Inches(6), Inches(0), Inches(4), Inches(7.5))
            title_size = Pt(max(24, 44 - (max(0, len(title_text) - 18) * 1.2)))
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(5.2), Inches(2.5))
            tf = tx.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = title_text
            p.font.name = theme.font_main; p.font.size = title_size; p.font.color.rgb = text_rgb; p.font.bold = True
            bt = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(5.2), Inches(4))
            btf = bt.text_frame; btf.word_wrap = True
            for b in bullets:
                bp = btf.add_paragraph(); bp.text = "• " + str(b)
                bp.font.name = theme.font_body; bp.font.size = Pt(18); bp.font.color.rgb = body_rgb

    prs.save(output_path)
    print(f"Success: Agnostic-Core v3 rendered at {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 4: sys.exit(1)
    create_slide_builder(sys.argv[1], sys.argv[2], sys.argv[3])
