"""
template_analyzer.py — Parse a PPTX template into a list of TextSlots.

Each TextSlot describes one text frame that the merge engine will overwrite,
addressed by a string slot key (see template_traversal.py for the scheme:
"42" plain shape, "42/17" group child, "42:r2c3" table cell). The same
traversal is used by the renderer, so addressing can never diverge.

All thresholds and limits are read from TemplateMergeConfig (sourced from
system_configs), never hardcoded.

Role inference heuristic (no LLM):
  - Placeholder type PP_PLACEHOLDER.TITLE / CENTER_TITLE → "title"
  - Placeholder type PP_PLACEHOLDER.SUBTITLE / BODY → "body"
  - Table cells → always "body" (the footnote/title position heuristics do
    not apply inside a table grid)
  - Non-placeholder in top `config.title_top_fraction` of slide height → "title"
  - Non-placeholder with area < `config.footnote_area_fraction` of slide → "footnote"
  - Everything else → "body"

Shape filtering (non-placeholder shapes only):
  - Empty text frame → skip
  - Existing text < config.shape_min_text_length → skip
  - Area > config.shape_bg_area_threshold of slide → skip (background)
  - Area < config.shape_min_area_threshold of slide → skip (decorative dot)
  - Table cells are exempt from the area filters (individual cells are
    legitimately tiny relative to the slide)
"""
import logging
from dataclasses import dataclass, field
from typing import List, Optional

from services.templates.template_config import TemplateMergeConfig
from services.templates.template_traversal import TextTarget, collect_text_targets

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False
    logger.warning("[TemplateAnalyzer] python-pptx not available.")


@dataclass
class TextSlot:
    slide_idx: int
    slot_key: str           # traversal key: "42" | "42/17" | "42:r2c3"
    shape_name: str
    role: str               # "title" | "body" | "footnote"
    char_limit: int
    hint: str               # existing text (topic hint for LLM)
    is_placeholder: bool
    placeholder_type: Optional[str] = None
    action: str = "rewrite" # "preserve" | "adapt" | "rewrite"
    kind: str = "shape"     # "shape" | "group_child" | "cell"


@dataclass
class SlideProfile:
    slide_idx: int
    slide_width_emu: int
    slide_height_emu: int
    slots: List[TextSlot] = field(default_factory=list)
    preserved_shapes: int = 0

    @property
    def hint(self) -> str:
        return " / ".join(s.hint for s in self.slots if s.role == "title" and s.hint)


def analyze_template(
    pptx_path: str,
    config: TemplateMergeConfig,
) -> List[SlideProfile]:
    """
    Open pptx_path and return one SlideProfile per slide.
    All filtering thresholds come from `config` (loaded from system_configs).
    """
    if not _PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is required for template analysis.")

    prs = Presentation(pptx_path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    slide_area = slide_w * slide_h

    profiles: List[SlideProfile] = []

    for slide_idx, slide in enumerate(prs.slides):
        profile = SlideProfile(
            slide_idx=slide_idx,
            slide_width_emu=int(slide_w),
            slide_height_emu=int(slide_h),
        )

        targets, preserved = collect_text_targets(slide, config.group_max_depth)
        profile.preserved_shapes = preserved

        for target in targets:
            try:
                slot = _build_slot(target, slide_idx, int(slide_h), int(slide_area), config)
                if slot is None:
                    profile.preserved_shapes += 1
                    continue
                profile.slots.append(slot)
            except Exception as exc:
                logger.warning(
                    f"[TemplateAnalyzer] slide {slide_idx} target "
                    f"'{target.name}' skipped: {exc}"
                )
                profile.preserved_shapes += 1

        profiles.append(profile)
        logger.info(
            f"[TemplateAnalyzer] slide {slide_idx}: {len(profile.slots)} text slots, "
            f"{profile.preserved_shapes} preserved shapes."
        )

    return profiles


# ─── private ──────────────────────────────────────────────────────────────────

def _build_slot(
    target: TextTarget,
    slide_idx: int,
    slide_height_emu: int,
    slide_area_emu: int,
    config: TemplateMergeConfig,
) -> Optional[TextSlot]:
    """Filter + classify one traversal target. Returns None when it must be preserved as-is."""
    existing_text = target.text_frame.text.strip()

    if not target.is_placeholder:
        # Skip non-placeholder targets without meaningful content
        if not existing_text:
            return None
        if len(existing_text) < config.shape_min_text_length:
            return None
        # Cells are exempt from slide-relative area filters
        if target.kind != "cell" and not _area_within_bounds(target, slide_area_emu, config):
            return None

    role = _infer_role(target, slide_height_emu, slide_area_emu, config)
    char_limit = _estimate_char_limit(target, role, existing_text, config)
    action = _infer_action(target.is_placeholder, role, existing_text, config)

    return TextSlot(
        slide_idx=slide_idx,
        slot_key=target.key,
        shape_name=target.name,
        role=role,
        char_limit=char_limit,
        hint=existing_text[:config.hint_max_chars],
        is_placeholder=target.is_placeholder,
        placeholder_type=_placeholder_type_str(target),
        action=action,
        kind=target.kind,
    )


def _area_within_bounds(
    target: TextTarget, slide_area_emu: int, config: TemplateMergeConfig
) -> bool:
    try:
        area = target.width * target.height
        if slide_area_emu <= 0:
            return True
        ratio = area / slide_area_emu
        if ratio > config.shape_bg_area_threshold:
            return False
        if ratio < config.shape_min_area_threshold:
            return False
    except Exception:
        pass
    return True


def _infer_role(
    target: TextTarget, slide_height_emu: int, slide_area_emu: int, config: TemplateMergeConfig
) -> str:
    if target.kind == "cell":
        return "body"

    if target.is_placeholder and target.shape is not None:
        try:
            from pptx.enum.shapes import PP_PLACEHOLDER
            ph_type = target.shape.placeholder_format.type
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                return "title"
            if ph_type in (PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY):
                return "body"
        except Exception:
            pass

    try:
        area = target.width * target.height

        if area < slide_area_emu * config.footnote_area_fraction:
            return "footnote"
        if target.top < slide_height_emu * config.title_top_fraction:
            return "title"
    except Exception:
        pass

    return "body"


def _estimate_char_limit(
    target: TextTarget, role: str, hint: str, config: TemplateMergeConfig
) -> int:
    hint_stripped = hint.strip()

    if role == "title":
        if hint_stripped and len(hint_stripped) <= config.short_hint_threshold:
            return max(len(hint_stripped) * config.short_hint_title_multiplier, 20)
        return config.title_char_limit

    if role == "footnote":
        return config.footnote_char_limit

    # Body
    if hint_stripped and len(hint_stripped) <= config.short_hint_threshold:
        return max(len(hint_stripped) * config.short_hint_body_multiplier, 30)

    try:
        w_in = target.width / 914400
        h_in = target.height / 914400
        estimated = int(w_in * h_in * config.chars_per_sq_inch)
        return max(config.body_char_limit_min, min(estimated, config.body_char_limit_max))
    except Exception:
        return config.body_char_limit_min


def _placeholder_type_str(target: TextTarget) -> Optional[str]:
    if not target.is_placeholder or target.shape is None:
        return None
    try:
        return str(target.shape.placeholder_format.type)
    except Exception:
        return None


def _infer_action(
    is_placeholder: bool,
    role: str,
    hint: str,
    config: TemplateMergeConfig,
) -> str:
    """
    Classify how the LLM should treat this slot:

      PRESERVE — do not touch; the existing text IS the correct content.
      ADAPT    — rewrite keeping the same semantic territory and approximate length.
      REWRITE  — free replacement from the knowledge base (default for placeholders).

    Priority order (first match wins):
      1. Footnotes are always preserved (legal/confidential text).
      2. Hints containing a preserve keyword are preserved regardless of length.
      3. Placeholder shapes (TITLE/BODY/SUBTITLE) are always rewritten.
      4. Non-placeholder with short hint → PRESERVE (structural label).
      5. Non-placeholder with medium hint → ADAPT (data to replace, structure to keep).
      6. Default → REWRITE.
    """
    # 1. Footnotes always preserved
    if role == "footnote":
        return "preserve"

    # 2. Legal / confidential keywords
    hint_lower = hint.lower()
    for kw in config.preserve_keywords.split(","):
        kw = kw.strip().lower()
        if kw and kw in hint_lower:
            return "preserve"

    # 3. Placeholder shapes → full rewrite
    if is_placeholder:
        return "rewrite"

    # 4. Short non-placeholder hint → structural label, preserve
    if len(hint) <= config.preserve_max_hint_chars:
        return "preserve"

    # 5. Medium non-placeholder hint → adapt (keep structure, replace data)
    if len(hint) <= config.adapt_max_hint_chars:
        return "adapt"

    # 6. Default
    return "rewrite"
