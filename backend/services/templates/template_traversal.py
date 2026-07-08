"""
template_traversal.py — Shared shape traversal for the Template Merge Engine.

Both the analyzer and the renderer walk slides through collect_text_targets()
so slot addressing can never diverge between analysis and render time.

Slot key scheme (string, stable within a slide):
  "42"       → top-level shape with shape_id 42
  "42/17"    → shape 17 inside group 42 (recursive: "42/17/9")
  "42:r2c3"  → cell at row 2, col 3 (0-indexed) of the table in shape 42

Coverage:
  - Plain shapes with a text frame → one target per shape.
  - GroupShapes → recursed up to max_group_depth; children keyed by path.
    Child geometry uses raw EMU values (group child coordinate space usually
    matches slide EMUs; role inference on children is best-effort).
  - Tables → one target per non-spanned cell; spanned cells (covered by a
    merge origin) are skipped so cell merges are never disturbed.
  - Charts, SmartArt, pictures, and anything else without an accessible
    text frame → counted in `preserved_count`, never yielded.
"""
import logging
from dataclasses import dataclass
from typing import List, Optional, Tuple

logger = logging.getLogger(__name__)

try:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


@dataclass
class TextTarget:
    """One replaceable text frame found while walking a slide."""
    key: str                 # slot key (see module docstring)
    text_frame: object       # pptx TextFrame
    name: str                # human-readable name for logs/report
    kind: str                # "shape" | "group_child" | "cell"
    is_placeholder: bool
    shape: object = None     # originating pptx shape (None-safe for cells)
    width: int = 0           # EMU (cells: column width × row height)
    height: int = 0          # EMU
    top: int = 0             # EMU (cells/group children: containing shape's top)


def collect_text_targets(slide, max_group_depth: int = 3) -> Tuple[List[TextTarget], int]:
    """
    Walk `slide` and return (targets, preserved_count) where `targets` is every
    text frame the merge engine could write to and `preserved_count` is the
    number of shapes/cells skipped because they carry no editable text frame
    (pictures, charts, groups beyond the depth cap, spanned cells, ...).
    """
    targets: List[TextTarget] = []
    preserved = 0
    for shape in slide.shapes:
        preserved += _walk(shape, "", 0, max_group_depth, targets)
    return targets, preserved


# ─── private ──────────────────────────────────────────────────────────────────

def _walk(shape, prefix: str, depth: int, max_depth: int, targets: List[TextTarget]) -> int:
    """Append targets found under `shape`; return count of preserved (non-target) items."""
    key = f"{prefix}{shape.shape_id}"
    preserved = 0

    shape_type = _safe_shape_type(shape)

    if shape_type == "group":
        if depth >= max_depth:
            logger.warning(
                f"[TemplateTraversal] group '{getattr(shape, 'name', '?')}' beyond "
                f"depth cap ({max_depth}) — children preserved as-is."
            )
            return 1
        for child in shape.shapes:
            preserved += _walk(child, f"{key}/", depth + 1, max_depth, targets)
        return preserved

    if shape_type == "table":
        return _walk_table(shape, key, targets)

    if getattr(shape, "has_text_frame", False):
        targets.append(TextTarget(
            key=key,
            text_frame=shape.text_frame,
            name=shape.name,
            kind="group_child" if depth > 0 else "shape",
            is_placeholder=bool(shape.is_placeholder) if depth == 0 else False,
            shape=shape,
            width=int(shape.width or 0),
            height=int(shape.height or 0),
            top=int(shape.top or 0),
        ))
        return 0

    return 1


def _walk_table(shape, key: str, targets: List[TextTarget]) -> int:
    preserved = 0
    try:
        table = shape.table
    except Exception:
        return 1

    top = int(shape.top or 0)
    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            try:
                if getattr(cell, "is_spanned", False):
                    preserved += 1
                    continue
                width = int(table.columns[c_idx].width or 0)
                height = int(row.height or 0)
                targets.append(TextTarget(
                    key=f"{key}:r{r_idx}c{c_idx}",
                    text_frame=cell.text_frame,
                    name=f"{shape.name}!r{r_idx}c{c_idx}",
                    kind="cell",
                    is_placeholder=False,
                    shape=None,
                    width=width,
                    height=height,
                    top=top,
                ))
            except Exception as exc:
                logger.warning(
                    f"[TemplateTraversal] table '{getattr(shape, 'name', '?')}' "
                    f"cell r{r_idx}c{c_idx} skipped: {exc}"
                )
                preserved += 1
    return preserved


def _safe_shape_type(shape) -> Optional[str]:
    """Return 'group' | 'table' | None without letting pptx quirks raise."""
    try:
        if _PPTX_AVAILABLE and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return "group"
    except Exception:
        pass
    try:
        if getattr(shape, "has_table", False):
            return "table"
    except Exception:
        pass
    return None
